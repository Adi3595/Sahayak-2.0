# main.py - COMPLETE VERSION WITH POSTGRESQL
import os
import io
import zipfile
import shutil
import uuid
import json
import random
import schedule
import time
import threading
import hashlib
from datetime import datetime, timedelta
from typing import List, Dict, Optional
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Request, BackgroundTasks, Depends, status
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from jose import JWTError, jwt
from passlib.context import CryptContext
from pydantic import BaseModel, EmailStr
from sqlalchemy.orm import Session
import base64
from io import BytesIO
import requests
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from groq import Groq

from utils import (
    save_upload,
    extract_text_from_file,
    generate_worksheets_with_groq,
    create_docx_from_text,
    create_pdf_from_html_optional,
    extract_topic_from_text,
    create_formatted_docx,
    create_formatted_pdf,
    format_grades_display
)

# Import database
from database import SessionLocal, engine, Base, User, UserSettings, UserProfile, Conversation, ConversationMessage, LessonPlan, GeneratedFile, StudentRecord, ProgressHistory, create_tables, test_connection

# ---------- Database Setup ----------
create_tables()

# Dependency to get DB session
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ---------- Authentication Setup ----------
SECRET_KEY = os.environ.get("SECRET_KEY")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/auth/login")

# ---------- Config ----------
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"
TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "templates")

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(TEMPLATE_DIR, exist_ok=True)

# ---------- Groq API ----------
import os

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
client = Groq(api_key=GROQ_API_KEY)

# ---------- App ----------
app = FastAPI(title="Sahayak - AI Teaching Assistant Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------- In-memory store for temporary jobs ----------
ppt_jobs = {}

# ---------- Pydantic Models ----------
class UserCreate(BaseModel):
    username: str
    email: EmailStr
    full_name: str
    password: str
    account_type: str = "teacher"

class UserLogin(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str
    user_id: int
    username: str
    full_name: str

class TokenData(BaseModel):
    username: Optional[str] = None

class StudentRecord(BaseModel):
    student_id: str
    name: str
    class_name: str
    grade: str
    progress: float
    status: str
    attendance: float
    last_updated: str

class ProgressUpdate(BaseModel):
    student_id: str
    class_name: str
    assignment_name: str
    score: float
    max_score: float
    assignment_type: str
    date_completed: str
    notes: Optional[str] = None

class StudentCreate(BaseModel):
    name: str
    class_name: str
    grade: str
    email: Optional[str] = None
    parent_contact: Optional[str] = None

class PPTRequest(BaseModel):
    topic: str
    slide_count: int
    style: str
    include_sections: List[str]
    subject: Optional[str] = None
    grade_level: Optional[str] = None

class DiagramRequest(BaseModel):
    topic: str
    style: str
    num_diagrams: int
    content_type: str
    text_content: Optional[str] = None

class SettingsUpdate(BaseModel):
    general: dict
    notifications: dict
    security: dict

class ProfileUpdate(BaseModel):
    personal: dict
    teaching: dict
    usage_stats: dict = None

class PasswordUpdate(BaseModel):
    current_password: str
    new_password: str
    confirm_password: str

class ExportDataRequest(BaseModel):
    data_type: str = "all"
    format: str = "json"

# ---------- Authentication Helpers ----------
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
        token_data = TokenData(username=username)
    except JWTError:
        raise credentials_exception
    
    user = db.query(User).filter(User.username == token_data.username).first()
    if user is None:
        raise credentials_exception
    if not user.is_active:
        raise HTTPException(status_code=400, detail="Inactive user")
    return user

async def get_current_active_user(current_user: User = Depends(get_current_user)):
    if not current_user.is_active:
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

# ---------- Authentication Endpoints ----------
@app.post("/api/auth/register", response_model=Token)
async def register(user_data: UserCreate, db: Session = Depends(get_db)):
    # Check if user exists
    existing_user = db.query(User).filter(
        (User.username == user_data.username) | (User.email == user_data.email)
    ).first()
    
    if existing_user:
        raise HTTPException(
            status_code=400,
            detail="Username or email already registered"
        )
    
    # Create new user
    hashed_password = get_password_hash(user_data.password)
    
    db_user = User(
        username=user_data.username,
        email=user_data.email,
        full_name=user_data.full_name,
        hashed_password=hashed_password,
        account_type=user_data.account_type,
        is_active=True,
        is_verified=True
    )
    
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    
    # Create default settings
    default_settings = initialize_default_settings()
    for settings_type, settings_data in default_settings.items():
        db_settings = UserSettings(
            user_id=db_user.id,
            settings_type=settings_type,
            settings_data=settings_data
        )
        db.add(db_settings)
    
    # Create default profile
    default_profile = initialize_default_profile()
    default_profile["personal"]["full_name"] = user_data.full_name
    default_profile["personal"]["email"] = user_data.email
    
    for profile_type, profile_data in default_profile.items():
        db_profile = UserProfile(
            user_id=db_user.id,
            profile_type=profile_type,
            profile_data=profile_data
        )
        db.add(db_profile)
    
    db.commit()
    
    # Create access token
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user_data.username}, expires_delta=access_token_expires
    )
    
    return Token(
        access_token=access_token,
        token_type="bearer",
        user_id=db_user.id,
        username=db_user.username,
        full_name=db_user.full_name
    )

@app.post("/api/auth/login", response_model=Token)
async def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.username == form_data.username).first()
    
    if not user or not verify_password(form_data.password, user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    if not user.is_active:
        raise HTTPException(status_code=400, detail="Inactive user")
    
    # Update last login
    user.last_login = datetime.utcnow()
    db.commit()
    
    # Create access token
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user.username}, expires_delta=access_token_expires
    )
    
    return Token(
        access_token=access_token,
        token_type="bearer",
        user_id=user.id,
        username=user.username,
        full_name=user.full_name
    )

@app.get("/api/auth/me")
async def read_users_me(current_user: User = Depends(get_current_active_user)):
    return {
        "id": current_user.id,
        "username": current_user.username,
        "email": current_user.email,
        "full_name": current_user.full_name,
        "account_type": current_user.account_type,
        "is_active": current_user.is_active,
        "created_at": current_user.created_at.isoformat(),
        "last_login": current_user.last_login.isoformat() if current_user.last_login else None
    }

# ---------- Helper Functions for Settings/Profile ----------
def initialize_default_settings():
    """Initialize default user settings"""
    return {
        "general": {
            "dark_mode": True,
            "email_notifications": False,
            "ai_suggestions": True,
            "language": "English",
            "timezone": "UTC",
            "date_format": "MM/DD/YYYY",
            "theme": "dark",
            "font_size": "medium",
            "auto_save": True
        },
        "notifications": {
            "lesson_reminders": True,
            "feature_updates": True,
            "weekly_reports": False,
            "assignment_alerts": True,
            "system_notifications": True,
            "email_notifications": False,
            "push_notifications": False
        },
        "security": {
            "two_factor_auth": False,
            "session_timeout": 30,
            "login_alerts": True,
            "data_backup": True,
            "privacy_mode": False
        },
        "privacy": {
            "data_collection": True,
            "analytics": True,
            "personalized_ads": False,
            "share_usage_data": True
        }
    }

def initialize_default_profile():
    """Initialize default user profile"""
    return {
        "personal": {
            "full_name": "",
            "email": "",
            "phone": "",
            "bio": "Experienced educator passionate about innovative teaching methods.",
            "avatar": "T",
            "join_date": datetime.now().strftime("%Y-%m-%d"),
            "location": "",
            "timezone": "UTC",
            "website": "",
            "social_links": {}
        },
        "teaching": {
            "primary_subject": "Mathematics",
            "secondary_subjects": ["Science", "English"],
            "grade_levels": ["Grade 1–3"],
            "languages": ["English"],
            "teaching_style": "Interactive and student-centered",
            "years_experience": 5,
            "specialization": "STEM Education",
            "certifications": ["Teaching License"],
            "teaching_philosophy": "",
            "classroom_management": ""
        },
        "preferences": {
            "theme": "dark",
            "font_size": "medium",
            "dashboard_layout": "standard",
            "default_view": "materials",
            "auto_save": True,
            "offline_mode": False,
            "accessibility": {
                "high_contrast": False,
                "text_to_speech": False,
                "keyboard_shortcuts": True
            }
        },
        "usage_stats": {
            "lessons_created": 0,
            "activities_generated": 0,
            "worksheets_created": 0,
            "presentations_made": 0,
            "diagrams_created": 0,
            "languages_used": 1,
            "total_students": 0,
            "plan": "Free",
            "active_since": datetime.now().strftime("%Y-%m-%d"),
            "last_active": datetime.now().isoformat()
        },
        "account": {
            "account_type": "teacher",
            "subscription": "free",
            "storage_used": "0 MB",
            "storage_limit": "1 GB",
            "last_login": datetime.now().isoformat(),
            "account_created": datetime.now().isoformat(),
            "account_status": "active",
            "billing_cycle": "none"
        }
    }

# ---------- Serve frontend pages ----------
@app.get("/", response_class=HTMLResponse)
async def get_home():
    file_path = os.path.join(os.path.dirname(__file__), "index.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h1>Welcome to Sahayak AI Teaching Assistant</h1>")

@app.get("/index.html")
async def serve_index():
    file_path = os.path.join(os.path.dirname(__file__), "index.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Dashboard page not found</h2>", status_code=404)

@app.get("/lesson_planning.html", response_class=HTMLResponse)
async def serve_lesson_plan():
    file_path = os.path.join(os.path.dirname(__file__), "lesson_planning.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Lesson Plan page not found</h2>", status_code=404)

@app.get("/Differentiated_materials.html", response_class=HTMLResponse)
async def serve_differentiated_materials():
    file_path = os.path.join(os.path.dirname(__file__), "Differentiated_materials.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Differentiated Materials page not found</h2>", status_code=404)

@app.get("/fun_activities.html", response_class=HTMLResponse)
async def serve_fun_activities():
    file_path = os.path.join(os.path.dirname(__file__), "fun_activities.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Fun Activities page not found</h2>", status_code=404)

@app.get("/blackboard_diagrams.html", response_class=HTMLResponse)
async def serve_blackboard_diagrams():
    file_path = os.path.join(os.path.dirname(__file__), "blackboard_diagrams.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Blackboard Diagrams page not found</h2>", status_code=404)

@app.get("/Student_tracking.html", response_class=HTMLResponse)
async def serve_student_progress():
    file_path = os.path.join(os.path.dirname(__file__), "Student_tracking.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Student Progress page not found</h2>", status_code=404)

@app.get("/knowledge_base.html", response_class=HTMLResponse)
async def serve_knowledge_base():
    file_path = os.path.join(os.path.dirname(__file__), "knowledge_base.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Knowledge Base page not found</h2>", status_code=404)

@app.get("/settings.html", response_class=HTMLResponse)
async def serve_settings_page():
    file_path = os.path.join(os.path.dirname(__file__), "settings.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Settings page not found</h2>", status_code=404)

@app.get("/profile.html", response_class=HTMLResponse)
async def serve_profile_page():
    file_path = os.path.join(os.path.dirname(__file__), "profile.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Profile page not found</h2>", status_code=404)

@app.get("/help.html", response_class=HTMLResponse)
async def serve_help_page():
    file_path = os.path.join(os.path.dirname(__file__), "help.html")
    if os.path.exists(file_path):
        return FileResponse(file_path)
    return HTMLResponse("<h2>Help page not found</h2>", status_code=404)

# ---------- Settings Endpoints ----------
@app.get("/api/settings")
async def get_settings(current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    """Get user settings from database"""
    try:
        settings = db.query(UserSettings).filter(
            UserSettings.user_id == current_user.id
        ).all()
        
        settings_dict = {}
        for setting in settings:
            settings_dict[setting.settings_type] = setting.settings_data
        
        # Ensure all settings types exist
        default_settings = initialize_default_settings()
        for setting_type in default_settings.keys():
            if setting_type not in settings_dict:
                # Create missing setting
                new_setting = UserSettings(
                    user_id=current_user.id,
                    settings_type=setting_type,
                    settings_data=default_settings[setting_type]
                )
                db.add(new_setting)
                settings_dict[setting_type] = default_settings[setting_type]
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "settings": settings_dict,
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching settings: {str(e)}")

@app.post("/api/settings/save")
async def save_settings(settings_data: SettingsUpdate, current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    """Save user settings to database"""
    try:
        # Update or create each settings type
        for settings_type, settings_value in settings_data.dict().items():
            existing = db.query(UserSettings).filter(
                UserSettings.user_id == current_user.id,
                UserSettings.settings_type == settings_type
            ).first()
            
            if existing:
                existing.settings_data = settings_value
                existing.updated_at = datetime.utcnow()
            else:
                new_setting = UserSettings(
                    user_id=current_user.id,
                    settings_type=settings_type,
                    settings_data=settings_value
                )
                db.add(new_setting)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Settings saved successfully",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error saving settings: {str(e)}")

@app.post("/api/settings/reset")
async def reset_settings(section: str = "all", current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    """Reset settings to default"""
    try:
        default_settings = initialize_default_settings()
        
        if section == "all":
            # Delete all settings for user
            db.query(UserSettings).filter(
                UserSettings.user_id == current_user.id
            ).delete()
            
            # Create default settings
            for settings_type, settings_data in default_settings.items():
                new_setting = UserSettings(
                    user_id=current_user.id,
                    settings_type=settings_type,
                    settings_data=settings_data
                )
                db.add(new_setting)
        
        elif section in default_settings:
            # Update specific section
            existing = db.query(UserSettings).filter(
                UserSettings.user_id == current_user.id,
                UserSettings.settings_type == section
            ).first()
            
            if existing:
                existing.settings_data = default_settings[section]
                existing.updated_at = datetime.utcnow()
            else:
                new_setting = UserSettings(
                    user_id=current_user.id,
                    settings_type=section,
                    settings_data=default_settings[section]
                )
                db.add(new_setting)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": f"Settings for '{section}' reset to default",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error resetting settings: {str(e)}")

# ---------- Profile Endpoints ----------
@app.get("/api/profile")
async def get_profile(current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    """Get user profile from database"""
    try:
        profiles = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id
        ).all()
        
        profile_dict = {}
        for profile in profiles:
            profile_dict[profile.profile_type] = profile.profile_data
        
        # Ensure all profile types exist
        default_profile = initialize_default_profile()
        for profile_type in default_profile.keys():
            if profile_type not in profile_dict:
                # Create missing profile
                new_profile = UserProfile(
                    user_id=current_user.id,
                    profile_type=profile_type,
                    profile_data=default_profile[profile_type]
                )
                db.add(new_profile)
                profile_dict[profile_type] = default_profile[profile_type]
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "profile": profile_dict,
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching profile: {str(e)}")

@app.post("/api/profile/save")
async def save_profile(profile_data: ProfileUpdate, current_user: User = Depends(get_current_active_user), db: Session = Depends(get_db)):
    """Save user profile to database"""
    try:
        # Update or create each profile type
        updates = profile_data.dict(exclude_unset=True)
        
        for profile_type, profile_value in updates.items():
            if profile_value is None:
                continue
                
            existing = db.query(UserProfile).filter(
                UserProfile.user_id == current_user.id,
                UserProfile.profile_type == profile_type
            ).first()
            
            if existing:
                # Merge existing data with new data
                if isinstance(profile_value, dict) and isinstance(existing.profile_data, dict):
                    existing.profile_data.update(profile_value)
                else:
                    existing.profile_data = profile_value
                existing.updated_at = datetime.utcnow()
            else:
                new_profile = UserProfile(
                    user_id=current_user.id,
                    profile_type=profile_type,
                    profile_data=profile_value
                )
                db.add(new_profile)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Profile saved successfully",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error saving profile: {str(e)}")

@app.post("/api/profile/update-password")
async def update_password(
    password_data: PasswordUpdate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Update user password in database"""
    try:
        # Verify current password
        if not verify_password(password_data.current_password, current_user.hashed_password):
            return JSONResponse(content={
                "status": "error",
                "message": "Current password is incorrect"
            }, status_code=400)
        
        # Validate new password
        if password_data.new_password != password_data.confirm_password:
            return JSONResponse(content={
                "status": "error",
                "message": "New passwords do not match"
            }, status_code=400)
        
        if len(password_data.new_password) < 8:
            return JSONResponse(content={
                "status": "error",
                "message": "Password must be at least 8 characters"
            }, status_code=400)
        
        # Update password
        current_user.hashed_password = get_password_hash(password_data.new_password)
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Password updated successfully",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error updating password: {str(e)}")

@app.get("/api/export-data")
async def export_data(
    data_type: str = "all",
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Export user data"""
    try:
        # Get settings
        settings = db.query(UserSettings).filter(
            UserSettings.user_id == current_user.id
        ).all()
        
        settings_dict = {}
        for setting in settings:
            settings_dict[setting.settings_type] = setting.settings_data
        
        # Get profile
        profiles = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id
        ).all()
        
        profile_dict = {}
        for profile in profiles:
            profile_dict[profile.profile_type] = profile.profile_data
        
        # Get conversations
        conversations = db.query(Conversation).filter(
            Conversation.user_id == current_user.id
        ).all()
        
        # Get lesson plans
        lesson_plans = db.query(LessonPlan).filter(
            LessonPlan.user_id == current_user.id
        ).all()
        
        # Get generated files
        generated_files = db.query(GeneratedFile).filter(
            GeneratedFile.user_id == current_user.id
        ).all()
        
        # Get students
        students = db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id
        ).all()
        
        # Get progress history
        progress_history = []
        for student in students:
            history = db.query(ProgressHistory).filter(
                ProgressHistory.student_id == student.student_id
            ).all()
            for entry in history:
                progress_history.append({
                    "student_id": entry.student_id,
                    "assignment_name": entry.assignment_name,
                    "score": entry.score,
                    "max_score": entry.max_score,
                    "percentage": entry.percentage,
                    "date_completed": entry.date_completed.isoformat()
                })
        
        # Prepare export data
        export_data = {
            "user_id": current_user.id,
            "username": current_user.username,
            "email": current_user.email,
            "full_name": current_user.full_name,
            "exported_at": datetime.now().isoformat(),
            "application": "Sahayak AI Teaching Assistant",
            "version": "1.0.0"
        }
        
        if data_type in ["all", "profile"]:
            export_data["profile"] = profile_dict
        
        if data_type in ["all", "settings"]:
            export_data["settings"] = settings_dict
        
        if data_type in ["all", "conversations"]:
            export_data["conversations"] = [
                {
                    "id": conv.conversation_id,
                    "title": conv.title,
                    "created_at": conv.created_at.isoformat(),
                    "updated_at": conv.updated_at.isoformat()
                }
                for conv in conversations
            ]
        
        if data_type in ["all", "lesson_plans"]:
            export_data["lesson_plans"] = [
                {
                    "id": lp.lesson_id,
                    "title": lp.title,
                    "topic": lp.topic,
                    "grade": lp.grade,
                    "created_at": lp.created_at.isoformat()
                }
                for lp in lesson_plans
            ]
        
        if data_type in ["all", "generated_files"]:
            export_data["generated_files"] = [
                {
                    "id": gf.file_id,
                    "filename": gf.filename,
                    "file_type": gf.file_type,
                    "created_at": gf.created_at.isoformat(),
                    "metadata": gf.metadata
                }
                for gf in generated_files
            ]
        
        if data_type in ["all", "students"]:
            export_data["students"] = [
                {
                    "id": student.student_id,
                    "name": student.name,
                    "class": student.class_name,
                    "grade": student.grade,
                    "progress": student.progress,
                    "status": student.status
                }
                for student in students
            ]
        
        if data_type in ["all", "progress"]:
            export_data["progress_history"] = progress_history
        
        # Create export file
        job_id = str(uuid.uuid4())
        filename = f"sahayak_export_{data_type}_{job_id[:8]}.json"
        filepath = os.path.join(OUTPUT_FOLDER, filename)
        
        with open(filepath, 'w') as f:
            json.dump(export_data, f, indent=2)
        
        # Track in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=filename,
            file_type="json",
            file_path=filepath,
            file_size=os.path.getsize(filepath),
            metadata={
                "export_type": data_type,
                "exported_at": datetime.now().isoformat()
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Data exported successfully",
            "filename": filename,
            "download_url": f"/download/{filename}",
            "data_type": data_type,
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error exporting data: {str(e)}")

@app.delete("/api/account")
async def delete_account(
    confirmation: str = None,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Delete user account"""
    try:
        if confirmation != "DELETE_MY_ACCOUNT":
            return JSONResponse(content={
                "status": "error",
                "message": "Confirmation phrase is required"
            }, status_code=400)
        
        # Mark user as inactive (soft delete)
        current_user.is_active = False
        current_user.username = f"deleted_{current_user.username}_{uuid.uuid4().hex[:8]}"
        current_user.email = f"deleted_{current_user.email}_{uuid.uuid4().hex[:8]}"
        
        # Mark all user data as inactive
        db.query(Conversation).filter(
            Conversation.user_id == current_user.id
        ).update({"is_active": False})
        
        db.query(LessonPlan).filter(
            LessonPlan.user_id == current_user.id
        ).update({"is_active": False})
        
        db.query(GeneratedFile).filter(
            GeneratedFile.user_id == current_user.id
        ).update({"is_active": False})
        
        db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id
        ).update({"is_active": False})
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Account deleted successfully",
            "timestamp": datetime.now().isoformat()
        })
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error deleting account: {str(e)}")

# ---------- Lesson Plan Generator ----------
@app.post("/generate/")
async def generate(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    try:
        data = await request.json()
        prompt = data.get("prompt", "")
        grade = data.get("grade")
        conversation_id = data.get("conversation_id") or str(uuid.uuid4())

        # Check if conversation exists in database
        conversation = db.query(Conversation).filter(
            Conversation.conversation_id == conversation_id,
            Conversation.user_id == current_user.id
        ).first()
        
        if not conversation:
            # Create new conversation in database
            conversation = Conversation(
                conversation_id=conversation_id,
                user_id=current_user.id,
                title=prompt[:50] + "..." if len(prompt) > 50 else prompt,
                created_at=datetime.utcnow(),
                updated_at=datetime.utcnow(),
                is_active=True
            )
            db.add(conversation)
            db.commit()
            db.refresh(conversation)

        # Save user message to database
        user_message = ConversationMessage(
            conversation_id=conversation_id,
            user_id=current_user.id,
            role="user",
            content=prompt,
            message_type="user_input",
            metadata={"grade": grade} if grade else {},
            created_at=datetime.utcnow()
        )
        db.add(user_message)
        db.commit()

        # Update conversation title
        conversation.updated_at = datetime.utcnow()
        if len(prompt) > 50 and conversation.title == prompt[:50] + "...":
            conversation.title = prompt[:50] + "..."
        db.commit()

        is_lesson_request = any(keyword in prompt.lower() for keyword in [
            "lesson plan", "lesson", "teaching plan", "class plan",
            "curriculum", "syllabus", "teaching", "educate", "teach",
            "classroom activity", "learning objectives", "lesson outline"
        ])

        if is_lesson_request:
            if not grade:
                bot_response = {
                    "response": "I'd be happy to create a lesson plan for you! For which grade level would you like me to design this lesson?",
                    "type": "grade_request",
                    "conversation_id": conversation_id,
                    "requires_grade": True
                }
                
                # Save assistant message to database
                assistant_message = ConversationMessage(
                    conversation_id=conversation_id,
                    user_id=current_user.id,
                    role="assistant",
                    content=bot_response["response"],
                    message_type=bot_response["type"],
                    metadata={"requires_grade": True},
                    created_at=datetime.utcnow()
                )
                db.add(assistant_message)
                db.commit()
                
                return JSONResponse(content=bot_response)

            topic = prompt.lower()
            for phrase in ["lesson plan", "lesson", "create", "make", "generate", "for grade", "about"]:
                topic = topic.replace(phrase, "")
            topic = topic.strip().capitalize() or "General Topic"

            lesson_prompt = f"""
You are an expert educator. Create a comprehensive, practical, and engaging lesson plan for teachers.

GRADE LEVEL: {grade}
TOPIC: {topic}
DURATION: 45-60 minutes

Format response with the following sections:
🎯 LESSON OVERVIEW
📚 LEARNING OBJECTIVES
🧩 MATERIALS & RESOURCES
⏰ LESSON PROCEDURE
🎨 DIFFERENTIATION STRATEGIES
📝 ASSESSMENT & FEEDBACK
💡 TEACHER TIPS

Keep it:
- Teacher-friendly
- Age-appropriate for grade {grade}
- Clear, engaging, and visually structured with emojis
- Practical and actionable
"""

            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": lesson_prompt}],
                temperature=0.7,
                max_tokens=2000
            )

            lesson_content = getattr(resp.choices[0].message, "content", str(resp))
            lesson_id = str(uuid.uuid4())
            lesson_plan_data = {
                "id": lesson_id,
                "topic": topic,
                "grade": grade,
                "content": lesson_content,
                "created_at": datetime.now().isoformat(),
                "conversation_id": conversation_id
            }

            # Save lesson plan to database
            lesson_plan = LessonPlan(
                lesson_id=lesson_id,
                user_id=current_user.id,
                conversation_id=conversation_id,
                title=topic,
                topic=topic,
                grade=grade,
                content=lesson_content,
                created_at=datetime.utcnow(),
                updated_at=datetime.utcnow(),
                is_active=True
            )
            db.add(lesson_plan)

            bot_response = {
                "response": lesson_content,
                "type": "lesson_plan",
                "grade": grade,
                "topic": topic,
                "conversation_id": conversation_id,
                "lesson_id": lesson_id
            }
        else:
            # Get recent conversation history from database
            recent_messages = db.query(ConversationMessage).filter(
                ConversationMessage.conversation_id == conversation_id,
                ConversationMessage.user_id == current_user.id
            ).order_by(ConversationMessage.created_at.desc()).limit(6).all()
            
            recent_messages.reverse()  # Oldest first
            
            conversation_history = []
            for msg in recent_messages:
                conversation_history.append({
                    "role": msg.role,
                    "content": msg.content
                })
            
            messages_for_api = [
                {"role": "system", "content": """
You are a helpful AI teaching assistant. 
Support teachers by providing actionable classroom ideas, 
teaching methods, and age-appropriate advice. 
If they ask for a lesson plan, ask their grade level first.
Be practical, encouraging, and focus on real classroom applications.
"""}
            ]
            
            for msg in conversation_history:
                messages_for_api.append({"role": msg["role"], "content": msg["content"]})

            messages_for_api.append({"role": "user", "content": prompt})

            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages_for_api,
                temperature=0.7,
                max_tokens=1000
            )

            chat_content = getattr(resp.choices[0].message, "content", str(resp))
            bot_response = {
                "response": chat_content,
                "type": "chat",
                "conversation_id": conversation_id
            }

        # Save assistant message to database
        assistant_message = ConversationMessage(
            conversation_id=conversation_id,
            user_id=current_user.id,
            role="assistant",
            content=bot_response["response"],
            message_type=bot_response["type"],
            metadata={
                "grade": grade if is_lesson_request else None,
                "topic": topic if is_lesson_request else None,
                "lesson_id": lesson_id if is_lesson_request else None
            } if is_lesson_request else {},
            created_at=datetime.utcnow()
        )
        db.add(assistant_message)
        
        # Update usage stats
        if is_lesson_request:
            profile = db.query(UserProfile).filter(
                UserProfile.user_id == current_user.id,
                UserProfile.profile_type == "usage_stats"
            ).first()
            if profile:
                stats = profile.profile_data
                stats["lessons_created"] = stats.get("lessons_created", 0) + 1
                stats["last_active"] = datetime.now().isoformat()
                profile.profile_data = stats
                profile.updated_at = datetime.utcnow()
        
        db.commit()

        print(f"\n✅ [{bot_response['type'].upper()}] Response generated for conversation {conversation_id} at {datetime.now()}")
        return JSONResponse(content=bot_response)

    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error generating response: {str(e)}")

# ---------- PDF Summarizer Endpoint ----------
@app.post("/summarize-pdf/")
async def summarize_pdf(
    file: UploadFile = File(...),
    summary_length: str = Form("medium"),
    page_start: int = Form(1),
    page_end: int = Form(None),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Dedicated endpoint for PDF summarization with page range support
    """
    try:
        # Save uploaded file
        saved_path = await save_upload(file, UPLOAD_FOLDER)
        if not saved_path:
            raise HTTPException(status_code=400, detail="Failed to save uploaded file")
        
        # Extract text from file with page range
        source_text = extract_text_from_file(saved_path, page_start=page_start, page_end=page_end)
        if not source_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the specified pages")
        
        # Determine summary length parameters
        length_params = {
            "short": {"max_tokens": 500, "temperature": 0.7},
            "medium": {"max_tokens": 1000, "temperature": 0.7},
            "detailed": {"max_tokens": 2000, "temperature": 0.7}
        }
        
        params = length_params.get(summary_length, length_params["medium"])
        
        # Create summarization prompt with page range info
        page_range_info = f"Pages {page_start}" + (f"-{page_end}" if page_end else "+")
        summarization_prompt = f"""
Please provide a comprehensive and well-structured summary of the following document content from {page_range_info}.
Focus on the main ideas, key concepts, and important details.

DOCUMENT CONTENT:
{source_text[:12000]}  # Limit context to avoid token limits

Please structure your summary with:
- Main topic and purpose
- Key points and concepts
- Important findings or conclusions
- Relevance and applications

Make the summary {summary_length} in length, focusing on clarity and coherence.
Avoid repeating the same information multiple times.
Provide a natural flow from introduction to conclusion.
"""
        
        # Generate summary using Groq
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system", 
                    "content": "You are an expert at summarizing educational content. Create clear, concise, and well-structured summaries that capture the essence of the document while maintaining accuracy and readability. Avoid repetition and ensure the summary flows naturally from introduction to conclusion."
                },
                {
                    "role": "user",
                    "content": summarization_prompt
                }
            ],
            temperature=params["temperature"],
            max_tokens=params["max_tokens"],
            top_p=0.9
        )
        
        summary_content = response.choices[0].message.content
        
        # Create output file
        job_id = str(uuid.uuid4())
        pdf_filename = f"summary_{summary_length}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        create_pdf_from_html_optional(f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .section {{ margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <h1>Document Summary</h1>
            <div class="metadata">
                <strong>Original File:</strong> {file.filename}<br>
                <strong>Summary Length:</strong> {summary_length}<br>
                <strong>Pages:</strong> {page_range_info}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {summary_content.replace('\n', '<br>')}
            </div>
        </body>
        </html>
        """, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "summary",
                "summary_length": summary_length,
                "original_file": file.filename,
                "page_range": page_range_info
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Clean up uploaded file
        try:
            os.remove(saved_path)
        except:
            pass
        
        db.commit()
        
        return JSONResponse(content={
            "summary": summary_content,
            "download_url": f"/download/{pdf_filename}",
            "filename": pdf_filename,
            "page_range": page_range_info
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Summarization failed: {str(e)}")

# ---------- Worksheet Generator ----------
@app.post("/generate-worksheets/")
async def generate_worksheets(
    file: UploadFile = File(...),
    subject: str = Form(None),
    grades: str = Form(...),
    difficulty: str = Form("moderate"),
    question_types: str = Form("mcq,theory"),
    output_format: str = Form("pdf"),
    page_start: int = Form(1),
    page_end: int = Form(None),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate worksheets and track in database
    """
    try:
        # Save uploaded file
        saved_path = await save_upload(file, UPLOAD_FOLDER)
        if not saved_path:
            raise HTTPException(status_code=400, detail="Failed to save uploaded file")
        
        # Extract text from file with page range
        source_text = extract_text_from_file(saved_path, page_start=page_start, page_end=page_end)
        if not source_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the specified pages")
        
        # Extract topic from source text if subject not provided
        extracted_topic = extract_topic_from_text(source_text)
        actual_subject = subject or extracted_topic
        
        # Parse grades and format nicely
        grade_list = [g.strip() for g in grades.split(",") if g.strip()]
        if not grade_list:
            raise HTTPException(status_code=400, detail="No valid grades provided")
        
        # Parse question types
        question_type_list = [q.strip().lower() for q in question_types.split(",") if q.strip()]
        if not question_type_list:
            question_type_list = ["mcq", "theory"]
        
        grades_display = format_grades_display(grade_list)
        
        # Generate worksheets
        worksheets = generate_worksheets_with_groq(
            source_text=source_text,
            subject=actual_subject,
            grades=grade_list,
            difficulty=difficulty,
            question_types=question_type_list,
            groq_client=client
        )
        
        # Create download links for each grade
        download_links = {}
        page_range_info = f"Pages {page_start}" + (f"-{page_end}" if page_end else "+")
        
        for grade, content in worksheets.items():
            # Create safe filename
            safe_topic = "".join(c for c in actual_subject if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_topic = safe_topic.replace(' ', '_')[:30]
            
            # Include question types in filename
            qtypes_str = "_".join(question_type_list)
            filename_base = f"{safe_topic}_Grade_{grade}_{difficulty}_{qtypes_str}"
            
            pdf_filename = f"{filename_base}.pdf"
            pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
            
            docx_filename = f"{filename_base}.docx"
            docx_path = os.path.join(OUTPUT_FOLDER, docx_filename)
            
            # Create formatted files
            create_formatted_pdf(content, actual_subject, grade, difficulty, pdf_path, question_type_list)
            create_formatted_docx(content, actual_subject, grade, difficulty, docx_path, question_type_list)
            
            # Track files in database
            for file_info in [
                (pdf_filename, pdf_path, "pdf"),
                (docx_filename, docx_path, "docx")
            ]:
                filename, filepath, filetype = file_info
                if os.path.exists(filepath):
                    generated_file = GeneratedFile(
                        file_id=str(uuid.uuid4()),
                        user_id=current_user.id,
                        filename=filename,
                        file_type=filetype,
                        file_path=filepath,
                        file_size=os.path.getsize(filepath),
                        metadata={
                            "type": "worksheet",
                            "subject": actual_subject,
                            "grade": grade,
                            "difficulty": difficulty,
                            "question_types": question_type_list,
                            "page_range": page_range_info
                        },
                        created_at=datetime.utcnow(),
                        expires_at=datetime.utcnow() + timedelta(days=7),
                        is_active=True
                    )
                    db.add(generated_file)
            
            download_links[grade] = {
                'pdf': f"/download/{pdf_filename}",
                'docx': f"/download/{docx_filename}",
                'preview': content[:300] + "..." if len(content) > 300 else content
            }
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["worksheets_created"] = stats.get("worksheets_created", 0) + len(worksheets)
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        # Clean up uploaded file
        try:
            os.remove(saved_path)
        except:
            pass
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "topic": actual_subject,
            "grades_display": grades_display,
            "page_range": page_range_info,
            "question_types": question_type_list,
            "worksheets": download_links,
            "message": f"Generated {len(grade_list)} worksheet(s) for {grades_display} with {', '.join(question_type_list)} questions"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Worksheet generation failed: {str(e)}")

# ---------- Topic-based Worksheet Generator ----------
@app.post("/generate-from-topics/")
async def generate_from_topics(
    main_topic: str = Form(...),
    subject: str = Form(...),
    grades: str = Form(...),
    difficulty: str = Form("moderate"),
    question_types: str = Form("mcq,theory"),
    subtopics: str = Form(""),
    output_format: str = Form("pdf"),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate worksheets directly from topics without file upload
    """
    try:
        if not main_topic.strip():
            raise HTTPException(status_code=400, detail="Main topic is required")
        
        if not subject.strip():
            raise HTTPException(status_code=400, detail="Subject is required")
        
        grade_list = [g.strip() for g in grades.split(",") if g.strip()]
        if not grade_list:
            raise HTTPException(status_code=400, detail="At least one grade level is required")
        
        # Parse question types
        question_type_list = [q.strip().lower() for q in question_types.split(",") if q.strip()]
        if not question_type_list:
            question_type_list = ["mcq", "theory"]
        
        subtopic_list = [s.strip() for s in subtopics.split(",") if s.strip()] if subtopics else []
        
        # Create source text from topics
        source_text = f"""
MAIN TOPIC: {main_topic}
SUBJECT: {subject}
GRADE LEVELS: {', '.join(grade_list)}
DIFFICULTY: {difficulty}
QUESTION TYPES: {', '.join(question_type_list)}

Please create comprehensive educational worksheets focusing on {main_topic}.
"""
        
        if subtopic_list:
            source_text += f"\nSPECIFICALLY COVER THESE SUBTOPICS: {', '.join(subtopic_list)}"
        
        # Generate worksheets
        worksheets = generate_worksheets_with_groq(
            source_text=source_text,
            subject=subject,
            grades=grade_list,
            difficulty=difficulty,
            question_types=question_type_list,
            groq_client=client
        )
        
        # Create download links
        download_links = {}
        safe_topic = "".join(c for c in main_topic if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_topic = safe_topic.replace(' ', '_')[:30]
        
        for grade, content in worksheets.items():
            # Include question types in filename
            qtypes_str = "_".join(question_type_list)
            filename_base = f"{safe_topic}_Grade_{grade}_{difficulty}_{qtypes_str}"
            
            pdf_filename = f"{filename_base}.pdf"
            pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
            
            docx_filename = f"{filename_base}.docx"
            docx_path = os.path.join(OUTPUT_FOLDER, docx_filename)
            
            create_formatted_pdf(content, subject, grade, difficulty, pdf_path, question_type_list)
            create_formatted_docx(content, subject, grade, difficulty, docx_path, question_type_list)
            
            # Track files in database
            for file_info in [
                (pdf_filename, pdf_path, "pdf"),
                (docx_filename, docx_path, "docx")
            ]:
                filename, filepath, filetype = file_info
                if os.path.exists(filepath):
                    generated_file = GeneratedFile(
                        file_id=str(uuid.uuid4()),
                        user_id=current_user.id,
                        filename=filename,
                        file_type=filetype,
                        file_path=filepath,
                        file_size=os.path.getsize(filepath),
                        metadata={
                            "type": "worksheet",
                            "subject": subject,
                            "grade": grade,
                            "difficulty": difficulty,
                            "question_types": question_type_list,
                            "subtopics": subtopic_list
                        },
                        created_at=datetime.utcnow(),
                        expires_at=datetime.utcnow() + timedelta(days=7),
                        is_active=True
                    )
                    db.add(generated_file)
            
            download_links[grade] = {
                'pdf': f"/download/{pdf_filename}",
                'docx': f"/download/{docx_filename}",
                'preview': content[:300] + "..." if len(content) > 300 else content
            }
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["worksheets_created"] = stats.get("worksheets_created", 0) + len(worksheets)
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success", 
            "topic": main_topic,
            "subject": subject,
            "grades_display": format_grades_display(grade_list),
            "question_types": question_type_list,
            "worksheets": download_links,
            "subtopics": subtopic_list,
            "message": f"Successfully generated worksheets for {main_topic}"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Topic-based generation failed: {str(e)}")

# ---------- Subtopics Generator ----------
@app.post("/generate-subtopics/")
async def generate_subtopics(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate relevant subtopics for a main topic
    """
    try:
        data = await request.json()
        main_topic = data.get("main_topic", "").strip()
        
        if not main_topic:
            raise HTTPException(status_code=400, detail="Main topic is required")
        
        # Generate subtopics using Groq
        subtopic_prompt = f"""
Given the main topic "{main_topic}", generate 5-7 relevant subtopics that would be appropriate for educational worksheets.
Return only a comma-separated list of subtopics without any additional text, numbering, or explanations.

Example format: Algebra, Geometry, Calculus, Statistics, Trigonometry

Focus on subtopics that are:
- Educational and teachable
- Appropriate for worksheet creation
- Cover different aspects of the main topic
- Suitable for various grade levels
"""
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an educational expert. Generate relevant, educational subtopics suitable for worksheet creation. Return only a comma-separated list without any additional text."
                },
                {
                    "role": "user",
                    "content": subtopic_prompt
                }
            ],
            temperature=0.7,
            max_tokens=200
        )
        
        subtopics_text = response.choices[0].message.content.strip()
        # Clean up the response
        subtopics_text = subtopics_text.replace('.', '').replace('\n', ',')
        subtopics_list = [s.strip() for s in subtopics_text.split(',') if s.strip()]
        
        # Remove duplicates and limit to 7
        unique_subtopics = []
        for subtopic in subtopics_list:
            if subtopic not in unique_subtopics and len(unique_subtopics) < 7:
                unique_subtopics.append(subtopic)
        
        return JSONResponse(content={
            "main_topic": main_topic,
            "subtopics": unique_subtopics
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Failed to generate subtopics: {str(e)}")

# ---------- PPT Creator Endpoints ----------
@app.post("/create-ppt/")
async def create_ppt(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(None),
    topic: str = Form(None),
    subject: str = Form("General"),
    grade_level: str = Form("9-12"),
    slide_count: int = Form(15),
    style: str = Form("academic"),
    include_sections: str = Form("title,outline,content,summary"),
    ppt_style: str = Form("professional"),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Create PowerPoint presentations from topics or documents
    """
    try:
        job_id = str(uuid.uuid4())
        ppt_jobs[job_id] = {
            "status": "processing",
            "created_at": datetime.now().isoformat(),
            "progress": 0,
            "user_id": current_user.id
        }
        
        # Start background task
        background_tasks.add_task(
            generate_ppt_background,
            job_id,
            file,
            topic,
            subject,
            grade_level,
            slide_count,
            style,
            include_sections.split(','),
            ppt_style,
            current_user.id,
            db
        )
        
        return JSONResponse(content={
            "status": "processing",
            "job_id": job_id,
            "message": "PPT generation started in background",
            "check_url": f"/check-ppt-status/{job_id}"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPT creation failed: {str(e)}")

async def generate_ppt_background(
    job_id: str,
    file: UploadFile,
    topic: str,
    subject: str,
    grade_level: str,
    slide_count: int,
    style: str,
    include_sections: List[str],
    ppt_style: str,
    user_id: int,
    db: Session
):
    """
    Background task to generate PPT
    """
    try:
        ppt_jobs[job_id]["progress"] = 10
        
        # Extract content from file or use topic
        content_text = ""
        if file and file.filename:
            # Save and extract text from file
            saved_path = await save_upload(file, UPLOAD_FOLDER)
            if saved_path:
                content_text = extract_text_from_file(saved_path)
                ppt_jobs[job_id]["progress"] = 30
                try:
                    os.remove(saved_path)
                except:
                    pass
        else:
            # Generate content from topic
            if not topic:
                raise ValueError("Topic is required when no file is provided")
            
            prompt = f"""
Create comprehensive educational content about: {topic}

SUBJECT: {subject}
GRADE LEVEL: {grade_level}
STYLE: {style}
NUMBER OF SECTIONS: {slide_count}

Please provide well-structured content suitable for a presentation.
Include:
- Clear introduction
- Main concepts and explanations
- Examples and applications
- Summary and key takeaways
"""
            
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert educator creating presentation content. Provide clear, structured, and engaging content suitable for slides."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                temperature=0.7,
                max_tokens=3000
            )
            
            content_text = response.choices[0].message.content
            ppt_jobs[job_id]["progress"] = 40
        
        # Generate slide content
        ppt_jobs[job_id]["progress"] = 50
        
        slide_prompt = f"""
Based on the following content, create {slide_count} slides for a presentation.

CONTENT:
{content_text[:5000]}

Please organize the content into {slide_count} slides with:
1. Title slide
2. Outline/Agenda slide
3. Content slides (main body)
4. Summary/Conclusion slide

For each slide, provide:
- Slide title
- 3-5 bullet points (concise and clear)
- Optional: image suggestions

Format as JSON with slides array.
"""
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """You are creating structured presentation slides. 
                    Return valid JSON with this format:
                    {
                        "slides": [
                            {
                                "title": "Slide Title",
                                "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
                                "image_suggestion": "optional image description"
                            }
                        ]
                    }"""
                },
                {
                    "role": "user",
                    "content": slide_prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        ppt_jobs[job_id]["progress"] = 70
        
        # Parse response
        try:
            slide_data = json.loads(response.choices[0].message.content)
        except:
            # If not valid JSON, create simple structure
            lines = response.choices[0].message.content.split('\n')
            slides = []
            current_slide = None
            
            for line in lines:
                if line.strip().startswith('# ') or line.strip().startswith('## '):
                    if current_slide:
                        slides.append(current_slide)
                    title = line.strip('# ').strip()
                    current_slide = {
                        "title": title,
                        "content": []
                    }
                elif line.strip().startswith('-') or line.strip().startswith('*'):
                    if current_slide:
                        content = line.strip('-* ').strip()
                        current_slide["content"].append(content)
            
            if current_slide:
                slides.append(current_slide)
            
            slide_data = {"slides": slides[:slide_count]}
        
        # Create PowerPoint
        ppt_jobs[job_id]["progress"] = 80
        presentation = create_powerpoint(slide_data["slides"], topic or "Presentation", subject, ppt_style)
        
        # Save presentation
        ppt_filename = f"presentation_{job_id[:8]}.pptx"
        ppt_path = os.path.join(OUTPUT_FOLDER, ppt_filename)
        presentation.save(ppt_path)
        
        # Track in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=user_id,
            filename=ppt_filename,
            file_type="pptx",
            file_path=ppt_path,
            file_size=os.path.getsize(ppt_path),
            metadata={
                "type": "presentation",
                "topic": topic,
                "subject": subject,
                "grade_level": grade_level,
                "slide_count": len(slide_data["slides"]),
                "style": ppt_style
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == user_id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["presentations_made"] = stats.get("presentations_made", 0) + 1
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        ppt_jobs[job_id]["progress"] = 100
        ppt_jobs[job_id]["status"] = "completed"
        ppt_jobs[job_id]["download_url"] = f"/download/{ppt_filename}"
        ppt_jobs[job_id]["filename"] = ppt_filename
        ppt_jobs[job_id]["slides_count"] = len(slide_data["slides"])
        
    except Exception as e:
        ppt_jobs[job_id]["status"] = "failed"
        ppt_jobs[job_id]["error"] = str(e)
        print(f"PPT generation failed for job {job_id}: {e}")

def create_powerpoint(slides: List[Dict], topic: str, subject: str, style: str):
    """
    Create PowerPoint presentation from slide data
    """
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define styles
    styles = {
        "professional": {
            "title_color": RGBColor(0, 32, 96),  # Dark blue
            "text_color": RGBColor(0, 0, 0),     # Black
            "accent_color": RGBColor(0, 112, 192), # Blue
            "bg_color": RGBColor(255, 255, 255)  # White
        },
        "academic": {
            "title_color": RGBColor(128, 0, 0),  # Maroon
            "text_color": RGBColor(0, 0, 0),     # Black
            "accent_color": RGBColor(192, 80, 77), # Red
            "bg_color": RGBColor(255, 255, 255)  # White
        },
        "creative": {
            "title_color": RGBColor(46, 125, 50), # Green
            "text_color": RGBColor(33, 33, 33),   # Dark gray
            "accent_color": RGBColor(76, 175, 80), # Light green
            "bg_color": RGBColor(255, 255, 255)  # White
        },
        "minimal": {
            "title_color": RGBColor(66, 66, 66),  # Gray
            "text_color": RGBColor(97, 97, 97),   # Light gray
            "accent_color": RGBColor(158, 158, 158), # Gray
            "bg_color": RGBColor(255, 255, 255)  # White
        }
    }
    
    style_config = styles.get(style, styles["professional"])
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = topic or "Educational Presentation"
    subtitle.text = f"{subject}\nGenerated by Sahayak AI"
    
    # Set title color
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = style_config["title_color"]
            run.font.size = Pt(44)
            run.font.bold = True
    
    # Set subtitle color
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = style_config["text_color"]
            run.font.size = Pt(20)
    
    # Content Slides
    for slide_data in slides:
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = slide_data.get("title", "Slide")
        
        # Set title style
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = style_config["title_color"]
                run.font.size = Pt(32)
                run.font.bold = True
        
        # Add content
        content = slide_data.get("content", [])
        text_frame = content_shape.text_frame
        text_frame.clear()  # Clear default text
        
        for i, bullet in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = style_config["text_color"]
        
        # Add slide number
        txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.2), Inches(0.8), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(len(prs.slides))
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(12)
        p.font.color.rgb = style_config["accent_color"]
    
    return prs

@app.get("/check-ppt-status/{job_id}")
async def check_ppt_status(job_id: str, current_user: User = Depends(get_current_active_user)):
    """
    Check status of PPT generation job
    """
    if job_id not in ppt_jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    # Check if job belongs to current user
    if ppt_jobs[job_id].get("user_id") != current_user.id:
        raise HTTPException(status_code=403, detail="Access denied")
    
    return JSONResponse(content=ppt_jobs[job_id])

@app.post("/create-ppt-from-doc/")
async def create_ppt_from_doc(
    file: UploadFile = File(...),
    style: str = Form("academic"),
    slide_count: int = Form(15),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Create PPT directly from uploaded document
    """
    try:
        job_id = str(uuid.uuid4())
        
        # Save uploaded file
        saved_path = await save_upload(file, UPLOAD_FOLDER)
        if not saved_path:
            raise HTTPException(status_code=400, detail="Failed to save uploaded file")
        
        # Extract text from file
        content_text = extract_text_from_file(saved_path)
        if not content_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the document")
        
        # Extract topic from text
        topic = extract_topic_from_text(content_text)
        
        # Generate slide content
        slide_prompt = f"""
Convert this document content into {slide_count} presentation slides:

CONTENT:
{content_text[:4000]}

Create structured slides with:
- Clear titles
- Concise bullet points (3-5 per slide)
- Logical flow from introduction to conclusion
"""
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """You are converting document content into presentation slides. 
                    Format as JSON:
                    {
                        "slides": [
                            {
                                "title": "Slide Title",
                                "content": ["Point 1", "Point 2", "Point 3"]
                            }
                        ]
                    }"""
                },
                {
                    "role": "user",
                    "content": slide_prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        # Parse response
        try:
            slide_data = json.loads(response.choices[0].message.content)
        except:
            # Fallback: create slides from paragraphs
            paragraphs = [p.strip() for p in content_text.split('\n\n') if p.strip()]
            slides = []
            
            for i, para in enumerate(paragraphs[:slide_count]):
                lines = para.split('\n')[:5]
                slides.append({
                    "title": f"Slide {i+1}",
                    "content": [line.strip('-•* ').strip() for line in lines if line.strip()]
                })
            
            slide_data = {"slides": slides}
        
        # Create PowerPoint
        presentation = create_powerpoint(slide_data["slides"], topic, "Document", style)
        
        # Save presentation
        ppt_filename = f"presentation_{job_id[:8]}.pptx"
        ppt_path = os.path.join(OUTPUT_FOLDER, ppt_filename)
        presentation.save(ppt_path)
        
        # Track in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=ppt_filename,
            file_type="pptx",
            file_path=ppt_path,
            file_size=os.path.getsize(ppt_path),
            metadata={
                "type": "presentation",
                "topic": topic,
                "subject": "Document",
                "slide_count": len(slide_data["slides"]),
                "style": style
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["presentations_made"] = stats.get("presentations_made", 0) + 1
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        # Clean up
        try:
            os.remove(saved_path)
        except:
            pass
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "topic": topic,
            "filename": ppt_filename,
            "download_url": f"/download/{ppt_filename}",
            "slides_count": len(slide_data["slides"]),
            "message": f"Created {len(slide_data['slides'])} slides from document"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"PPT creation failed: {str(e)}")

@app.post("/create-ppt-from-topic/")
async def create_ppt_from_topic(
    topic: str = Form(...),
    subject: str = Form("General"),
    grade_level: str = Form("9-12"),
    slide_count: int = Form(15),
    style: str = Form("academic"),
    include_sections: str = Form("title,outline,content,summary"),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Create PPT directly from topic
    """
    try:
        if not topic.strip():
            raise HTTPException(status_code=400, detail="Topic is required")
        
        job_id = str(uuid.uuid4())
        
        # Generate content from topic
        prompt = f"""
Create comprehensive educational content for a presentation about: {topic}

SUBJECT: {subject}
GRADE LEVEL: {grade_level}
NUMBER OF SLIDES: {slide_count}
STYLE: {style}

Please provide content organized into {slide_count} slides covering:
1. Introduction to {topic}
2. Key concepts and principles
3. Examples and applications
4. Important facts and figures
5. Summary and conclusion

Make the content engaging and educational.
"""
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are creating educational presentation content. Provide well-structured, clear, and informative content suitable for slides."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2500
        )
        
        content_text = response.choices[0].message.content
        
        # Generate slides from content
        slide_prompt = f"""
Organize this content into exactly {slide_count} presentation slides:

CONTENT:
{content_text}

Create slides with:
- Clear, descriptive titles
- 3-5 concise bullet points per slide
- Logical progression of ideas
- Educational focus
"""
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are organizing content into presentation slides. Format as JSON with slides array."
                },
                {
                    "role": "user",
                    "content": slide_prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        # Parse slides
        try:
            slide_data = json.loads(response.choices[0].message.content)
        except:
            # Create simple slide structure
            lines = content_text.split('\n')
            slides = []
            current_title = None
            current_content = []
            
            for line in lines:
                if line.strip().startswith('# ') or line.strip().startswith('## ') or line.strip().startswith('### '):
                    if current_title and current_content:
                        slides.append({
                            "title": current_title,
                            "content": current_content[:5]  # Limit to 5 points
                        })
                    current_title = line.strip('# ').strip()
                    current_content = []
                elif line.strip().startswith('-') or line.strip().startswith('*'):
                    current_content.append(line.strip('-* ').strip())
                elif line.strip() and len(line.strip()) > 10:
                    current_content.append(line.strip())
            
            if current_title and current_content:
                slides.append({
                    "title": current_title,
                    "content": current_content[:5]
                })
            
            slide_data = {"slides": slides[:slide_count]}
        
        # Create PowerPoint
        presentation = create_powerpoint(slide_data["slides"], topic, subject, style)
        
        # Save presentation
        ppt_filename = f"presentation_{job_id[:8]}.pptx"
        ppt_path = os.path.join(OUTPUT_FOLDER, ppt_filename)
        presentation.save(ppt_path)
        
        # Track in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=ppt_filename,
            file_type="pptx",
            file_path=ppt_path,
            file_size=os.path.getsize(ppt_path),
            metadata={
                "type": "presentation",
                "topic": topic,
                "subject": subject,
                "grade_level": grade_level,
                "slide_count": len(slide_data["slides"]),
                "style": style
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["presentations_made"] = stats.get("presentations_made", 0) + 1
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "topic": topic,
            "filename": ppt_filename,
            "download_url": f"/download/{ppt_filename}",
            "slides_count": len(slide_data["slides"]),
            "message": f"Created presentation on '{topic}' with {len(slide_data['slides'])} slides"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"PPT creation failed: {str(e)}")

# ---------- Fun Learning Activities Generator ----------
@app.post("/generate-activity/")
async def generate_activity(
    subject: str = Form(...),
    topic: str = Form(...),
    grade_level: str = Form(...),
    activity_type: str = Form(None),
    duration: str = Form(None),
    resources: str = Form(""),
    learning_objective: str = Form(""),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate fun learning activities and games based on curriculum and resources
    """
    try:
        if not subject or not topic or not grade_level:
            raise HTTPException(status_code=400, detail="Subject, topic, and grade level are required")
        
        # Create activity generation prompt
        activity_prompt = f"""
Create a fun, engaging learning activity or game for teachers to use in the classroom.

SUBJECT: {subject}
TOPIC: {topic}
GRADE LEVEL: {grade_level}
LEARNING OBJECTIVE: {learning_objective or "Make learning engaging and memorable"}
AVAILABLE RESOURCES: {resources or "Standard classroom materials"}
DURATION: {duration or "Flexible"}
ACTIVITY TYPE: {activity_type or "Any engaging format"}

Please provide a comprehensive activity plan with:

1. ACTIVITY TITLE: Creative and engaging name
2. ACTIVITY TYPE: Game, group activity, individual task, etc.
3. DURATION: Estimated time required
4. MATERIALS NEEDED: List of required resources
5. STEP-BY-STEP INSTRUCTIONS: Clear, actionable steps for teachers
6. LEARNING OUTCOMES: What students will learn/practice
7. DIFFERENTIATION STRATEGIES: How to adapt for different learners
8. ASSESSMENT IDEAS: How to measure learning
9. VARIATIONS: Alternative ways to run the activity

Make the activity:
- Highly engaging and interactive
- Age-appropriate for the grade level
- Connected to the learning objective
- Practical with available resources
- Easy to implement for teachers
- Fun and memorable for students

Focus on active learning, student participation, and making the content come alive.
"""

        # Generate activity using Groq
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """You are an expert educational game designer and activity creator. 
                    Create fun, engaging, and educational activities that make learning enjoyable. 
                    Focus on interactive, hands-on experiences that promote active learning.
                    Provide practical, ready-to-use activities that teachers can implement easily."""
                },
                {
                    "role": "user",
                    "content": activity_prompt
                }
            ],
            temperature=0.8,
            max_tokens=3000,
            top_p=0.9
        )
        
        activity_content = response.choices[0].message.content
        
        # Create a PDF file for the activity
        job_id = str(uuid.uuid4())
        pdf_filename = f"activity_{topic.replace(' ', '_')}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        create_pdf_from_html_optional(f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .section {{ margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <h1>Fun Learning Activity</h1>
            <div class="metadata">
                <strong>Topic:</strong> {topic}<br>
                <strong>Subject:</strong> {subject}<br>
                <strong>Grade Level:</strong> {grade_level}<br>
                <strong>Activity Type:</strong> {activity_type or 'Mixed'}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {activity_content.replace('\n', '<br>')}
            </div>
        </body>
        </html>
        """, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "activity",
                "topic": topic,
                "subject": subject,
                "grade_level": grade_level,
                "activity_type": activity_type or "mixed"
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["activities_generated"] = stats.get("activities_generated", 0) + 1
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "activity": activity_content,
            "subject": subject,
            "topic": topic,
            "grade_level": grade_level,
            "activity_type": activity_type or "mixed",
            "download_url": f"/download/{pdf_filename}",
            "message": f"Created fun activity for {topic} in {subject}"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Activity generation failed: {str(e)}")

# ---------- Blackboard Diagrams Generator ----------
@app.post("/generate-visual-aids/")
async def generate_visual_aids(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate comprehensive visual aids with actual diagram images
    """
    try:
        data = await request.json()
        topic = data.get("topic", "")
        grade_level = data.get("grade_level", "middle school")
        lesson_duration = data.get("lesson_duration", 45)
        learning_objectives = data.get("learning_objectives", "")
        
        if not topic.strip():
            raise HTTPException(status_code=400, detail="Topic is required")
        
        # First, get the visual aids content
        prompt = f"""
Create a comprehensive set of visual aids for teaching: {topic}

GRADE LEVEL: {grade_level}
LESSON DURATION: {lesson_duration} minutes
LEARNING OBJECTIVES: {learning_objectives}

Please provide a complete visual teaching toolkit including:

1. BLACKBOARD DIAGRAMS (3-5 different types):
   - Main concept diagram
   - Process/flow diagram
   - Comparison chart
   - Timeline or sequence
   - Summary diagram

2. For EACH diagram, provide:
   - Clear title and purpose
   - Step-by-step drawing instructions
   - Simple shapes and elements to use
   - Labeling suggestions
   - Teaching tips

3. VISUAL TEACHING STRATEGIES

Focus on creating practical, classroom-ready visuals.
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """You are an expert educational consultant specializing in visual teaching methods. 
                    Create comprehensive visual aid packages with clear, implementable diagrams."""
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=3000
        )
        
        visual_aids_content = response.choices[0].message.content
        
        # Generate actual diagram images
        diagram_images = await generate_diagram_images(topic, grade_level)
        
        # Create a PDF file for the visual aids
        job_id = str(uuid.uuid4())
        pdf_filename = f"visual_aids_{topic.replace(' ', '_')}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        # Create HTML with embedded images
        html_content = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .diagram {{ margin: 20px 0; border: 1px solid #ddd; padding: 10px; }}
                img {{ max-width: 100%; height: auto; }}
            </style>
        </head>
        <body>
            <h1>Visual Teaching Aids</h1>
            <div class="metadata">
                <strong>Topic:</strong> {topic}<br>
                <strong>Grade Level:</strong> {grade_level}<br>
                <strong>Lesson Duration:</strong> {lesson_duration} minutes<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {visual_aids_content.replace('\n', '<br>')}
            </div>
        """
        
        # Add diagram images
        for i, diagram in enumerate(diagram_images):
            if diagram.get('image_data'):
                html_content += f"""
                <div class="diagram">
                    <h3>Diagram {i+1}: {diagram.get('title', 'Diagram')}</h3>
                    <img src="{diagram['image_data']}" alt="{diagram.get('title', 'Diagram')}">
                </div>
                """
        
        html_content += """
        </body>
        </html>
        """
        
        create_pdf_from_html_optional(html_content, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "visual_aids",
                "topic": topic,
                "grade_level": grade_level,
                "diagram_count": len(diagram_images)
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["diagrams_created"] = stats.get("diagrams_created", 0) + len(diagram_images)
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "topic": topic,
            "grade_level": grade_level,
            "lesson_duration": lesson_duration,
            "visual_aids": visual_aids_content,
            "diagrams": diagram_images,
            "download_url": f"/download/{pdf_filename}",
            "created_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Visual aids generation failed: {str(e)}")

async def generate_diagram_images(topic: str, grade_level: str):
    """Generate actual diagram images based on the topic"""
    diagrams = []
    
    try:
        # Generate different types of diagrams
        diagram_types = [
            "main_concept",
            "flowchart", 
            "comparison",
            "timeline",
            "summary"
        ]
        
        for diagram_type in diagram_types:
            image_data = await create_diagram_image(topic, diagram_type, grade_level)
            if image_data:
                diagrams.append({
                    "type": diagram_type,
                    "title": f"{topic} - {diagram_type.replace('_', ' ').title()}",
                    "image_data": image_data,
                    "filename": f"{topic}_{diagram_type}.png"
                })
        
        return diagrams
        
    except Exception as e:
        print(f"Diagram image generation failed: {e}")
        return []

async def create_diagram_image(topic: str, diagram_type: str, grade_level: str):
    """Create a specific diagram image using matplotlib"""
    try:
        # Create a figure with blackboard-like background
        fig, ax = plt.subplots(figsize=(10, 8))
        fig.patch.set_facecolor('#2D5016')  # Blackboard green
        ax.set_facecolor('#2D5016')
        
        # Set up the plot area
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 8)
        ax.axis('off')
        
        # Add title
        title_text = f"{topic}\n{diagram_type.replace('_', ' ').title()}"
        ax.text(5, 7.5, title_text, ha='center', va='center', 
                fontsize=14, color='white', weight='bold',
                bbox=dict(boxstyle="round,pad=0.3", facecolor='#5D4037', alpha=0.8))
        
        # Generate different diagram types
        if diagram_type == "main_concept":
            await draw_main_concept_diagram(ax, topic)
        elif diagram_type == "flowchart":
            await draw_flowchart_diagram(ax, topic)
        elif diagram_type == "comparison":
            await draw_comparison_diagram(ax, topic)
        elif diagram_type == "timeline":
            await draw_timeline_diagram(ax, topic)
        elif diagram_type == "summary":
            await draw_summary_diagram(ax, topic)
        
        # Add watermark
        ax.text(5, 0.3, "Sahayak Visual Aid", ha='center', va='center',
                fontsize=10, color='white', alpha=0.7, style='italic')
        
        # Convert to base64
        buffer = BytesIO()
        plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight', 
                   facecolor=fig.get_facecolor(), edgecolor='none')
        buffer.seek(0)
        
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        plt.close(fig)
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error creating {diagram_type} diagram: {e}")
        return None

async def draw_main_concept_diagram(ax, topic):
    """Draw a main concept diagram"""
    # Central concept circle
    center_circle = patches.Circle((5, 4), 1.5, fill=True, facecolor='white', 
                                 edgecolor='black', linewidth=2)
    ax.add_patch(center_circle)
    ax.text(5, 4, "Main\nConcept", ha='center', va='center', 
            fontsize=10, weight='bold')
    
    # Surrounding concepts
    concepts = ["Feature 1", "Feature 2", "Feature 3", "Feature 4"]
    angles = [45, 135, 225, 315]
    
    for i, (concept, angle) in enumerate(zip(concepts, angles)):
        rad = np.radians(angle)
        x = 5 + 3 * np.cos(rad)
        y = 4 + 3 * np.sin(rad)
        
        # Connection line
        ax.plot([5 + 1.5 * np.cos(rad), x - 0.5 * np.cos(rad)], 
                [4 + 1.5 * np.sin(rad), y - 0.5 * np.sin(rad)], 
                'w-', linewidth=2)
        
        # Concept box
        box = patches.Rectangle((x-1, y-0.5), 2, 1, fill=True, 
                              facecolor='lightblue', edgecolor='black')
        ax.add_patch(box)
        ax.text(x, y, concept, ha='center', va='center', fontsize=8)

async def draw_flowchart_diagram(ax, topic):
    """Draw a simple flowchart"""
    steps = [
        (5, 6.5, "Start", 'oval'),
        (5, 5.5, "Step 1", 'rectangle'),
        (5, 4.5, "Decision", 'diamond'),
        (3, 3.5, "Option A", 'rectangle'),
        (7, 3.5, "Option B", 'rectangle'),
        (5, 2.5, "End", 'oval')
    ]
    
    for x, y, text, shape in steps:
        if shape == 'oval':
            ellipse = patches.Ellipse((x, y), 2, 0.8, fill=True, 
                                    facecolor='lightgreen', edgecolor='black')
            ax.add_patch(ellipse)
        elif shape == 'rectangle':
            rect = patches.Rectangle((x-1, y-0.4), 2, 0.8, fill=True,
                                  facecolor='lightyellow', edgecolor='black')
            ax.add_patch(rect)
        elif shape == 'diamond':
            diamond = patches.Polygon([(x, y+0.4), (x+1, y), (x, y-0.4), (x-1, y)], 
                                    fill=True, facecolor='lightcoral', edgecolor='black')
            ax.add_patch(diamond)
        
        ax.text(x, y, text, ha='center', va='center', fontsize=8)
    
    # Draw connections
    connections = [
        [(5, 6.1), (5, 5.9)],
        [(5, 5.1), (5, 4.9)],
        [(4, 4.5), (3.5, 4.5), (3.5, 3.9)],
        [(6, 4.5), (6.5, 4.5), (6.5, 3.9)],
        [(3, 3.1), (3, 2.9), (5, 2.9)],
        [(7, 3.1), (7, 2.9), (5, 2.9)]
    ]
    
    for connection in connections:
        x_vals, y_vals = zip(*connection)
        ax.plot(x_vals, y_vals, 'k-', linewidth=2)

async def draw_comparison_diagram(ax, topic):
    """Draw a comparison chart"""
    categories = ["Category A", "Category B", "Category C"]
    values_a = [7, 5, 8]
    values_b = [4, 9, 6]
    
    x = np.arange(len(categories))
    width = 0.35
    
    # Create bars
    bars1 = ax.bar(x - width/2, values_a, width, label='Option 1', 
                  color='skyblue', edgecolor='black')
    bars2 = ax.bar(x + width/2, values_b, width, label='Option 2', 
                  color='lightcoral', edgecolor='black')
    
    # Customize appearance for blackboard
    ax.set_xticks(x)
    ax.set_xticklabels(categories, color='white')
    ax.tick_params(colors='white')
    ax.legend(facecolor='#5D4037', edgecolor='white', labelcolor='white')
    
    # Add value labels on bars
    for bar in bars1:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height}', ha='center', va='bottom', 
                color='white', weight='bold')
    
    for bar in bars2:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height}', ha='center', va='bottom', 
                color='white', weight='bold')

async def draw_timeline_diagram(ax, topic):
    """Draw a timeline diagram"""
    events = [
        (1, "Event 1", "Start"),
        (3, "Event 2", "Development"),
        (5, "Event 3", "Progress"),
        (7, "Event 4", "Milestone"),
        (9, "Event 5", "Completion")
    ]
    
    # Draw timeline
    ax.plot([0.5, 9.5], [4, 4], 'w-', linewidth=3)
    
    for x, event, phase in events:
        # Event marker
        ax.plot([x, x], [3.8, 4.2], 'w-', linewidth=2)
        
        # Event circle
        circle = patches.Circle((x, 5), 0.3, fill=True, 
                              facecolor='yellow', edgecolor='black')
        ax.add_patch(circle)
        ax.text(x, 5, event, ha='center', va='center', fontsize=8, weight='bold')
        
        # Phase label
        ax.text(x, 3.5, phase, ha='center', va='top', 
                fontsize=7, color='white', style='italic')

async def draw_summary_diagram(ax, topic):
    """Draw a summary/mind map diagram"""
    # Central topic
    center_circle = patches.Circle((5, 4), 0.8, fill=True, 
                                 facecolor='gold', edgecolor='black')
    ax.add_patch(center_circle)
    ax.text(5, 4, "Summary", ha='center', va='center', 
            fontsize=9, weight='bold')
    
    # Main points
    points = [
        (2, 6, "Key Point 1", 45),
        (8, 6, "Key Point 2", 135),
        (2, 2, "Key Point 3", 315),
        (8, 2, "Key Point 4", 225)
    ]
    
    for x, y, text, angle in points:
        # Connection line
        ax.plot([5, x], [4, y], 'w-', linewidth=2)
        
        # Point box
        box = patches.Rectangle((x-1.2, y-0.4), 2.4, 0.8, fill=True,
                              facecolor='lightgreen', edgecolor='black')
        ax.add_patch(box)
        ax.text(x, y, text, ha='center', va='center', fontsize=8)

# ---------- Enhanced Diagram Generator with Images ----------
@app.post("/generate-diagram/")
async def generate_diagram(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate blackboard diagram instructions with actual images
    """
    try:
        data = await request.json()
        topic = data.get("topic", "")
        style = data.get("style", "simple")
        num_diagrams = data.get("num_diagrams", 1)
        content_type = data.get("content_type", "general")
        text_content = data.get("text_content", "")
        
        if not topic.strip():
            raise HTTPException(status_code=400, detail="Topic is required")
        
        # Get diagram instructions
        prompt = f"Create {num_diagrams} simple blackboard diagrams for: {topic}"
        if content_type == "summary" and text_content:
            prompt += f"\nBased on this summary: {text_content}"
        elif content_type == "text_content" and text_content:
            prompt += f"\nBased on this content: {text_content}"
        
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educator who creates simple, clear blackboard diagrams."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        diagram_instructions = response.choices[0].message.content
        
        # Generate diagram images
        diagram_images = []
        for i in range(min(num_diagrams, 3)):  # Limit to 3 images max
            diagram_type = ["main_concept", "flowchart", "comparison"][i % 3]
            image_data = await create_diagram_image(topic, diagram_type, "middle school")
            if image_data:
                diagram_images.append({
                    "type": diagram_type,
                    "title": f"{topic} - Diagram {i+1}",
                    "image_data": image_data,
                    "filename": f"{topic}_diagram_{i+1}.png"
                })
        
        # Create a PDF file for the diagrams
        job_id = str(uuid.uuid4())
        pdf_filename = f"diagrams_{topic.replace(' ', '_')}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        # Create HTML with embedded images
        html_content = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .diagram {{ margin: 20px 0; border: 1px solid #ddd; padding: 10px; }}
                img {{ max-width: 100%; height: auto; }}
            </style>
        </head>
        <body>
            <h1>Blackboard Diagrams</h1>
            <div class="metadata">
                <strong>Topic:</strong> {topic}<br>
                <strong>Style:</strong> {style}<br>
                <strong>Number of Diagrams:</strong> {num_diagrams}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {diagram_instructions.replace('\n', '<br>')}
            </div>
        """
        
        # Add diagram images
        for i, diagram in enumerate(diagram_images):
            if diagram.get('image_data'):
                html_content += f"""
                <div class="diagram">
                    <h3>Diagram {i+1}: {diagram.get('title', 'Diagram')}</h3>
                    <img src="{diagram['image_data']}" alt="{diagram.get('title', 'Diagram')}">
                </div>
                """
        
        html_content += """
        </body>
        </html>
        """
        
        create_pdf_from_html_optional(html_content, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "diagrams",
                "topic": topic,
                "style": style,
                "num_diagrams": num_diagrams,
                "diagram_count": len(diagram_images)
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["diagrams_created"] = stats.get("diagrams_created", 0) + len(diagram_images)
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "topic": topic,
            "style": style,
            "num_diagrams": num_diagrams,
            "content_type": content_type,
            "instructions": diagram_instructions,
            "diagrams": diagram_images,
            "download_url": f"/download/{pdf_filename}",
            "created_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Diagram generation failed: {str(e)}")

# ---------- Download Diagram Images ----------
@app.post("/download-diagrams/")
async def download_diagrams(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Download generated diagrams as a ZIP file
    """
    try:
        data = await request.json()
        topic = data.get("topic", "")
        diagrams = data.get("diagrams", [])
        
        if not diagrams:
            raise HTTPException(status_code=400, detail="No diagrams to download")
        
        # Create a unique directory for this download
        job_id = str(uuid.uuid4())
        output_dir = os.path.join(OUTPUT_FOLDER, f"diagrams_{job_id}")
        os.makedirs(output_dir, exist_ok=True)
        
        # Save each diagram
        saved_files = []
        for diagram in diagrams:
            if diagram.get('image_data'):
                # Extract base64 data
                image_data = diagram['image_data'].split(',')[1]
                image_bytes = base64.b64decode(image_data)
                
                filename = diagram.get('filename', f"{topic}_diagram.png")
                file_path = os.path.join(output_dir, filename)
                
                with open(file_path, 'wb') as f:
                    f.write(image_bytes)
                saved_files.append(file_path)
        
        # Create ZIP file
        zip_filename = f"{topic}_diagrams.zip"
        zip_path = os.path.join(OUTPUT_FOLDER, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_path in saved_files:
                zipf.write(file_path, os.path.basename(file_path))
        
        # Track ZIP file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=zip_filename,
            file_type="zip",
            file_path=zip_path,
            file_size=os.path.getsize(zip_path),
            metadata={
                "type": "diagrams_zip",
                "topic": topic,
                "diagram_count": len(diagrams)
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        # Clean up individual files
        for file_path in saved_files:
            try:
                os.remove(file_path)
            except:
                pass
        try:
            os.rmdir(output_dir)
        except:
            pass
        
        db.commit()
        
        return JSONResponse(content={
            "download_url": f"/download/{zip_filename}",
            "filename": zip_filename,
            "diagram_count": len(diagrams),
            "message": f"Download ready for {len(diagrams)} diagrams"
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Download creation failed: {str(e)}")

# ---------- Knowledge Base Endpoints ----------
@app.post("/explain-concept/")
async def explain_concept(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Explain complex concepts with simple explanations, analogies, and local language support
    """
    try:
        data = await request.json()
        concept = data.get("concept", "")
        grade_level = data.get("grade_level", "general")
        language = data.get("language", "english")
        include_analogies = data.get("include_analogies", True)
        include_examples = data.get("include_examples", True)
        include_visuals = data.get("include_visuals", False)
        
        if not concept.strip():
            raise HTTPException(status_code=400, detail="Concept is required")
        
        # Create explanation prompt
        prompt = f"""
Explain the concept: "{concept}"

GRADE LEVEL: {grade_level}
LANGUAGE: {language}
INCLUDE ANALOGIES: {include_analogies}
INCLUDE EXAMPLES: {include_examples}
INCLUDE VISUALS: {include_visuals}

Please provide a comprehensive explanation with:

1. SIMPLE DEFINITION: Easy-to-understand definition
2. KEY POINTS: 3-5 main characteristics or components
3. {f"ANALOGIES: 2-3 relatable analogies from everyday life" if include_analogies else ""}
4. {f"EXAMPLES: Practical examples and applications" if include_examples else ""}
5. COMMON MISCONCEPTIONS: What people often get wrong
6. REAL-WORLD CONNECTIONS: How this applies to daily life
7. {f"VISUAL DESCRIPTION: Simple way to visualize this concept" if include_visuals else ""}
8. RELATED CONCEPTS: What to learn next

Make it:
- Age-appropriate for {grade_level}
- Clear and engaging
- Practical and relatable
- Culturally appropriate for {language} speakers
- Free of unnecessary jargon
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": f"""You are an expert educator who explains complex concepts in simple, engaging ways.
                    You specialize in making difficult topics accessible to all age groups.
                    You provide relatable analogies and practical examples.
                    You adapt your explanations to be culturally relevant."""
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=3000,
            top_p=0.9
        )
        
        explanation_content = response.choices[0].message.content
        
        # Generate related concepts
        related_concepts = await generate_related_concepts(concept, grade_level)
        
        # Create a PDF file for the explanation
        job_id = str(uuid.uuid4())
        pdf_filename = f"explanation_{concept.replace(' ', '_')}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        create_pdf_from_html_optional(f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .section {{ margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <h1>Concept Explanation</h1>
            <div class="metadata">
                <strong>Concept:</strong> {concept}<br>
                <strong>Grade Level:</strong> {grade_level}<br>
                <strong>Language:</strong> {language}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {explanation_content.replace('\n', '<br>')}
            </div>
        </body>
        </html>
        """, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "explanation",
                "concept": concept,
                "grade_level": grade_level,
                "language": language
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "concept": concept,
            "grade_level": grade_level,
            "language": language,
            "explanation": explanation_content,
            "related_concepts": related_concepts,
            "download_url": f"/download/{pdf_filename}",
            "created_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Concept explanation failed: {str(e)}")

async def generate_related_concepts(concept: str, grade_level: str):
    """Generate related concepts for further learning"""
    try:
        prompt = f"""
Given the concept "{concept}" for {grade_level} level, suggest 5-7 related concepts that would help build deeper understanding.
Return only a comma-separated list without additional text.

Focus on:
- Prerequisite concepts
- Complementary topics
- Advanced extensions
- Practical applications
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "Generate relevant, educational concepts that build on the main topic. Return only a comma-separated list."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=200
        )
        
        concepts_text = response.choices[0].message.content.strip()
        concepts_list = [c.strip() for c in concepts_text.split(',') if c.strip()]
        
        return concepts_list[:7]  # Limit to 7 concepts
        
    except Exception as e:
        print(f"Related concepts generation failed: {e}")
        return []

@app.post("/compare-concepts/")
async def compare_concepts(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Compare two or more concepts with similarities and differences
    """
    try:
        data = await request.json()
        concepts = data.get("concepts", [])
        comparison_focus = data.get("focus", "similarities and differences")
        grade_level = data.get("grade_level", "general")
        
        if len(concepts) < 2:
            raise HTTPException(status_code=400, detail="At least two concepts are required for comparison")
        
        concepts_text = ", ".join(concepts)
        
        prompt = f"""
Compare these concepts: {concepts_text}

COMPARISON FOCUS: {comparison_focus}
GRADE LEVEL: {grade_level}

Please provide a comprehensive comparison with:

1. BASIC DEFINITIONS: Simple definition of each concept
2. KEY SIMILARITIES: What they have in common
3. KEY DIFFERENCES: How they are distinct
4. RELATIONSHIPS: How they connect or interact
5. WHEN TO USE EACH: Practical applications and contexts
6. COMMON CONFUSIONS: What people mix up
7. LEARNING TIPS: How to remember the differences

Make the comparison:
- Clear and structured
- Practical and relevant
- Age-appropriate
- Helpful for understanding
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at comparing and contrasting concepts. You highlight both similarities and differences in clear, educational ways."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2500
        )
        
        comparison_content = response.choices[0].message.content
        
        # Create a PDF file for the comparison
        job_id = str(uuid.uuid4())
        pdf_filename = f"comparison_{'_'.join([c[:10] for c in concepts])}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        create_pdf_from_html_optional(f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .section {{ margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <h1>Concept Comparison</h1>
            <div class="metadata">
                <strong>Concepts:</strong> {concepts_text}<br>
                <strong>Focus:</strong> {comparison_focus}<br>
                <strong>Grade Level:</strong> {grade_level}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {comparison_content.replace('\n', '<br>')}
            </div>
        </body>
        </html>
        """, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "comparison",
                "concepts": concepts,
                "comparison_focus": comparison_focus,
                "grade_level": grade_level
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "concepts": concepts,
            "comparison_focus": comparison_focus,
            "grade_level": grade_level,
            "comparison": comparison_content,
            "download_url": f"/download/{pdf_filename}",
            "created_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Concept comparison failed: {str(e)}")

@app.post("/generate-analogies/")
async def generate_analogies(
    request: Request,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate multiple analogies for a concept
    """
    try:
        data = await request.json()
        concept = data.get("concept", "")
        context = data.get("context", "general")
        num_analogies = data.get("num_analogies", 3)
        grade_level = data.get("grade_level", "general")
        
        if not concept.strip():
            raise HTTPException(status_code=400, detail="Concept is required")
        
        prompt = f"""
Generate {num_analogies} creative and relatable analogies for: "{concept}"

CONTEXT: {context}
GRADE LEVEL: {grade_level}

For each analogy, provide:
1. The analogy itself
2. How it relates to the concept
3. Why it helps understanding
4. Potential limitations

Focus on analogies that are:
- Easy to understand
- Culturally relevant
- Age-appropriate
- Memorable and engaging
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are a master of creating perfect analogies that make complex concepts instantly understandable."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.8,
            max_tokens=2000
        )
        
        analogies_content = response.choices[0].message.content
        
        # Create a PDF file for the analogies
        job_id = str(uuid.uuid4())
        pdf_filename = f"analogies_{concept.replace(' ', '_')}_{job_id[:8]}.pdf"
        pdf_path = os.path.join(OUTPUT_FOLDER, pdf_filename)
        
        create_pdf_from_html_optional(f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #333; border-bottom: 2px solid #6C63FF; padding-bottom: 10px; }}
                .metadata {{ background: #f5f5f5; padding: 15px; border-radius: 5px; margin-bottom: 20px; }}
                .content {{ margin-top: 20px; }}
                .section {{ margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <h1>Analogies for: {concept}</h1>
            <div class="metadata">
                <strong>Concept:</strong> {concept}<br>
                <strong>Context:</strong> {context}<br>
                <strong>Number of Analogies:</strong> {num_analogies}<br>
                <strong>Grade Level:</strong> {grade_level}<br>
                <strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M')}
            </div>
            <div class="content">
                {analogies_content.replace('\n', '<br>')}
            </div>
        </body>
        </html>
        """, pdf_path)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=pdf_filename,
            file_type="pdf",
            file_path=pdf_path,
            file_size=os.path.getsize(pdf_path) if os.path.exists(pdf_path) else 0,
            metadata={
                "type": "analogies",
                "concept": concept,
                "context": context,
                "num_analogies": num_analogies,
                "grade_level": grade_level
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "concept": concept,
            "context": context,
            "grade_level": grade_level,
            "analogies": analogies_content,
            "download_url": f"/download/{pdf_filename}",
            "created_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Analogy generation failed: {str(e)}")

# ---------- Student Progress Tracking Endpoints ----------
@app.get("/api/students/")
async def get_all_students(
    class_filter: str = None,
    status_filter: str = None,
    grade_filter: str = None,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Get all students with optional filtering
    """
    try:
        query = db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id,
            StudentRecord.is_active == True
        )
        
        # Apply filters
        if class_filter and class_filter != "all":
            query = query.filter(StudentRecord.class_name == class_filter)
        if status_filter and status_filter != "all":
            query = query.filter(StudentRecord.status == status_filter)
        if grade_filter and grade_filter != "all":
            query = query.filter(StudentRecord.grade == grade_filter)
        
        students = query.order_by(StudentRecord.last_updated.desc()).all()
        
        students_dict = {}
        for student in students:
            students_dict[student.student_id] = {
                "student_id": student.student_id,
                "name": student.name,
                "class_name": student.class_name,
                "grade": student.grade,
                "email": student.email,
                "parent_contact": student.parent_contact,
                "progress": student.progress,
                "attendance": student.attendance,
                "status": student.status,
                "last_updated": student.last_updated.isoformat()
            }
        
        return JSONResponse(content={
            "students": students_dict,
            "total_count": len(students),
            "filters_applied": {
                "class": class_filter,
                "status": status_filter,
                "grade": grade_filter
            }
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching students: {str(e)}")

@app.post("/api/students/")
async def create_student(
    student_data: StudentCreate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Create a new student record in database
    """
    try:
        student_id = str(uuid.uuid4())
        
        # Create student record
        student_record = StudentRecord(
            student_id=student_id,
            user_id=current_user.id,
            name=student_data.name,
            class_name=student_data.class_name,
            grade=student_data.grade,
            email=student_data.email,
            parent_contact=student_data.parent_contact,
            progress=0.0,
            attendance=100.0,
            status="average",
            created_at=datetime.utcnow(),
            last_updated=datetime.utcnow(),
            is_active=True
        )
        
        db.add(student_record)
        
        # Update usage stats
        profile = db.query(UserProfile).filter(
            UserProfile.user_id == current_user.id,
            UserProfile.profile_type == "usage_stats"
        ).first()
        if profile:
            stats = profile.profile_data
            stats["total_students"] = stats.get("total_students", 0) + 1
            stats["last_active"] = datetime.now().isoformat()
            profile.profile_data = stats
            profile.updated_at = datetime.utcnow()
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "student_id": student_id,
            "message": f"Student {student_data.name} created successfully"
        })
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error creating student: {str(e)}")

@app.get("/api/students/{student_id}")
async def get_student_details(
    student_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Get detailed information for a specific student
    """
    try:
        student = db.query(StudentRecord).filter(
            StudentRecord.student_id == student_id,
            StudentRecord.user_id == current_user.id,
            StudentRecord.is_active == True
        ).first()
        
        if not student:
            raise HTTPException(status_code=404, detail="Student not found or access denied")
        
        # Get progress history
        history = db.query(ProgressHistory).filter(
            ProgressHistory.student_id == student_id
        ).order_by(ProgressHistory.recorded_at.desc()).all()
        
        history_list = []
        for entry in history:
            history_list.append({
                "entry_id": entry.entry_id,
                "assignment_name": entry.assignment_name,
                "assignment_type": entry.assignment_type,
                "score": entry.score,
                "max_score": entry.max_score,
                "percentage": entry.percentage,
                "date_completed": entry.date_completed.isoformat(),
                "notes": entry.notes,
                "recorded_at": entry.recorded_at.isoformat()
            })
        
        # Calculate statistics
        recent_scores = [entry.percentage for entry in history[:10]]  # Last 10 assignments
        average_score = sum(recent_scores) / len(recent_scores) if recent_scores else 0
        
        return JSONResponse(content={
            "student": {
                "student_id": student.student_id,
                "name": student.name,
                "class_name": student.class_name,
                "grade": student.grade,
                "email": student.email,
                "parent_contact": student.parent_contact,
                "progress": student.progress,
                "attendance": student.attendance,
                "status": student.status,
                "created_at": student.created_at.isoformat(),
                "last_updated": student.last_updated.isoformat()
            },
            "progress_history": history_list,
            "statistics": {
                "average_score": round(average_score, 2),
                "total_assignments": len(history),
                "recent_trend": calculate_trend_db(history),
                "attendance_rate": student.attendance
            }
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching student details: {str(e)}")

@app.post("/api/progress/")
async def update_student_progress(
    progress_data: ProgressUpdate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Update student progress in database
    """
    try:
        # Check if student belongs to current user
        student = db.query(StudentRecord).filter(
            StudentRecord.student_id == progress_data.student_id,
            StudentRecord.user_id == current_user.id
        ).first()
        
        if not student:
            raise HTTPException(status_code=404, detail="Student not found or access denied")
        
        # Calculate percentage
        percentage = (progress_data.score / progress_data.max_score) * 100
        
        # Create progress entry
        entry_id = str(uuid.uuid4())
        progress_entry = ProgressHistory(
            entry_id=entry_id,
            student_id=progress_data.student_id,
            assignment_name=progress_data.assignment_name,
            assignment_type=progress_data.assignment_type,
            score=progress_data.score,
            max_score=progress_data.max_score,
            percentage=round(percentage, 2),
            date_completed=datetime.fromisoformat(progress_data.date_completed.replace('Z', '+00:00')),
            notes=progress_data.notes,
            recorded_at=datetime.utcnow()
        )
        
        db.add(progress_entry)
        
        # Update student's overall progress
        await update_student_overall_progress_db(progress_data.student_id, current_user.id, db)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": "Progress updated successfully",
            "progress_entry": {
                "entry_id": entry_id,
                "student_id": progress_data.student_id,
                "assignment_name": progress_data.assignment_name,
                "percentage": round(percentage, 2)
            }
        })
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error updating progress: {str(e)}")

def calculate_trend_db(history: List[ProgressHistory]) -> str:
    """
    Calculate performance trend (improving, declining, stable)
    """
    if len(history) < 2:
        return "stable"
    
    recent = history[:5]  # Last 5 entries
    if len(recent) < 2:
        return "stable"
    
    first_avg = sum(entry.percentage for entry in recent[:2]) / 2
    last_avg = sum(entry.percentage for entry in recent[-2:]) / 2
    
    if last_avg - first_avg > 5:
        return "improving"
    elif first_avg - last_avg > 5:
        return "declining"
    else:
        return "stable"

async def update_student_overall_progress_db(student_id: str, user_id: int, db: Session):
    """Update student's overall progress based on recent performance"""
    try:
        # Get recent progress entries
        recent_entries = db.query(ProgressHistory).filter(
            ProgressHistory.student_id == student_id
        ).order_by(ProgressHistory.recorded_at.desc()).limit(10).all()
        
        if not recent_entries:
            return
        
        # Calculate weighted average
        total_weight = 0
        weighted_sum = 0
        
        for i, entry in enumerate(recent_entries):
            weight = i + 1  # More recent = higher weight
            total_weight += weight
            weighted_sum += entry.percentage * weight
        
        new_progress = weighted_sum / total_weight if total_weight > 0 else 0
        
        # Update student record
        student = db.query(StudentRecord).filter(
            StudentRecord.student_id == student_id,
            StudentRecord.user_id == user_id
        ).first()
        
        if student:
            student.progress = round(new_progress, 2)
            student.last_updated = datetime.utcnow()
            
            # Update status based on progress
            if new_progress >= 90:
                student.status = "excellent"
            elif new_progress >= 80:
                student.status = "good"
            elif new_progress >= 70:
                student.status = "average"
            else:
                student.status = "poor"
    
    except Exception as e:
        print(f"Error updating overall progress for {student_id}: {e}")

@app.get("/api/progress/overview")
async def get_progress_overview(
    class_name: str = None,
    time_range: str = "current_semester",
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Get overview statistics and charts data
    """
    try:
        # Filter students by class
        query = db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id,
            StudentRecord.is_active == True
        )
        
        if class_name and class_name != "all":
            query = query.filter(StudentRecord.class_name == class_name)
        
        filtered_students = query.all()
        
        # Calculate statistics
        total_students = len(filtered_students)
        if total_students == 0:
            return JSONResponse(content={
                "total_students": 0,
                "average_progress": 0,
                "attendance_rate": 0,
                "grade_distribution": {},
                "performance_trends": []
            })
        
        # Calculate averages
        total_progress = sum(student.progress for student in filtered_students)
        total_attendance = sum(student.attendance for student in filtered_students)
        
        average_progress = total_progress / total_students
        average_attendance = total_attendance / total_students
        
        # Calculate grade distribution
        grade_distribution = {
            "excellent": 0,  # 90-100%
            "good": 0,       # 80-89%
            "average": 0,    # 70-79%
            "poor": 0        # <70%
        }
        
        for student in filtered_students:
            progress = student.progress
            if progress >= 90:
                grade_distribution["excellent"] += 1
            elif progress >= 80:
                grade_distribution["good"] += 1
            elif progress >= 70:
                grade_distribution["average"] += 1
            else:
                grade_distribution["poor"] += 1
        
        # Generate performance trends (last 8 weeks)
        performance_trends = generate_performance_trends_db(filtered_students, time_range, db)
        
        return JSONResponse(content={
            "total_students": total_students,
            "average_progress": round(average_progress, 2),
            "attendance_rate": round(average_attendance, 2),
            "grade_distribution": grade_distribution,
            "performance_trends": performance_trends
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating overview: {str(e)}")

def generate_performance_trends_db(students: List[StudentRecord], time_range: str, db: Session):
    """Generate performance trends data for charts"""
    trends = []
    weeks = 8
    
    for week in range(weeks):
        # Simulate data (in real app, query actual historical data)
        week_data = {
            "week": f"Week {week + 1}",
            "average_score": 70 + (week * 2) + (random.random() * 10 - 5),  # Simulated improvement
            "top_performer": 95 + (random.random() * 5 - 2.5),
            "class_average": 75 + (week * 1.5) + (random.random() * 8 - 4),
            "assignments_completed": len(students) * (0.8 + (week * 0.05))  # Simulated increase
        }
        trends.append(week_data)
    
    return trends

@app.post("/api/progress/bulk-update")
async def bulk_update_progress(
    updates: List[ProgressUpdate],
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Update progress for multiple students at once
    """
    try:
        results = []
        for update in updates:
            try:
                # Check if student belongs to current user
                student = db.query(StudentRecord).filter(
                    StudentRecord.student_id == update.student_id,
                    StudentRecord.user_id == current_user.id
                ).first()
                
                if not student:
                    results.append({
                        "student_id": update.student_id,
                        "status": "error",
                        "message": "Student not found or access denied"
                    })
                    continue
                
                # Calculate percentage
                percentage = (update.score / update.max_score) * 100
                
                # Create progress entry
                entry_id = str(uuid.uuid4())
                progress_entry = ProgressHistory(
                    entry_id=entry_id,
                    student_id=update.student_id,
                    assignment_name=update.assignment_name,
                    assignment_type=update.assignment_type,
                    score=update.score,
                    max_score=update.max_score,
                    percentage=round(percentage, 2),
                    date_completed=datetime.fromisoformat(update.date_completed.replace('Z', '+00:00')),
                    notes=update.notes,
                    recorded_at=datetime.utcnow()
                )
                
                db.add(progress_entry)
                
                # Update student's overall progress
                await update_student_overall_progress_db(update.student_id, current_user.id, db)
                
                results.append({
                    "student_id": update.student_id,
                    "status": "success",
                    "message": "Progress updated"
                })
                
            except Exception as e:
                results.append({
                    "student_id": update.student_id,
                    "status": "error",
                    "message": str(e)
                })
        
        db.commit()
        
        return JSONResponse(content={
            "results": results,
            "total_processed": len(updates),
            "successful": len([r for r in results if r['status'] == 'success'])
        })
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error in bulk update: {str(e)}")

@app.post("/api/reports/generate")
async def generate_progress_report(
    report_type: str = Form("class_summary"),
    class_name: str = Form(None),
    student_ids: str = Form(None),
    include_charts: bool = Form(True),
    format: str = Form("pdf"),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Generate comprehensive progress reports
    """
    try:
        # Parse student IDs if provided
        target_student_ids = []
        if student_ids:
            target_student_ids = [sid.strip() for sid in student_ids.split(",") if sid.strip()]
        
        # Get data based on report type
        if report_type == "class_summary":
            report_data = await generate_class_summary_report(class_name, include_charts, current_user.id, db)
        elif report_type == "individual_reports":
            report_data = await generate_individual_reports(target_student_ids, include_charts, current_user.id, db)
        elif report_type == "comparison_report":
            report_data = await generate_comparison_report(target_student_ids, include_charts, current_user.id, db)
        else:
            raise HTTPException(status_code=400, detail="Invalid report type")
        
        # Generate file
        job_id = str(uuid.uuid4())
        filename = f"progress_report_{report_type}_{job_id[:8]}.{format}"
        filepath = os.path.join(OUTPUT_FOLDER, filename)
        
        if format == "pdf":
            create_progress_report_pdf(report_data, filepath)
        else:
            create_progress_report_docx(report_data, filepath)
        
        # Track file in database
        generated_file = GeneratedFile(
            file_id=str(uuid.uuid4()),
            user_id=current_user.id,
            filename=filename,
            file_type=format,
            file_path=filepath,
            file_size=os.path.getsize(filepath),
            metadata={
                "type": "progress_report",
                "report_type": report_type,
                "class_name": class_name,
                "student_count": len(target_student_ids) if target_student_ids else 0
            },
            created_at=datetime.utcnow(),
            expires_at=datetime.utcnow() + timedelta(days=7),
            is_active=True
        )
        db.add(generated_file)
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "download_url": f"/download/{filename}",
            "filename": filename,
            "report_type": report_type,
            "generated_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error generating report: {str(e)}")

async def generate_class_summary_report(class_name: str, include_charts: bool, user_id: int, db: Session):
    """Generate class summary report data"""
    # Get students
    query = db.query(StudentRecord).filter(
        StudentRecord.user_id == user_id,
        StudentRecord.is_active == True
    )
    
    if class_name and class_name != "all":
        query = query.filter(StudentRecord.class_name == class_name)
    
    students = query.all()
    
    # Calculate statistics
    total_students = len(students)
    if total_students > 0:
        avg_progress = sum(s.progress for s in students) / total_students
        avg_attendance = sum(s.attendance for s in students) / total_students
    else:
        avg_progress = 0
        avg_attendance = 0
    
    return {
        "report_type": "class_summary",
        "class_name": class_name or "All Classes",
        "generated_at": datetime.now().isoformat(),
        "summary": {
            "total_students": total_students,
            "average_progress": round(avg_progress, 2),
            "attendance_rate": round(avg_attendance, 2)
        },
        "students": [
            {
                "name": s.name,
                "class_name": s.class_name,
                "grade": s.grade,
                "progress": s.progress,
                "status": s.status,
                "last_updated": s.last_updated.isoformat()
            }
            for s in students
        ],
        "include_charts": include_charts
    }

async def generate_individual_reports(student_ids: List[str], include_charts: bool, user_id: int, db: Session):
    """Generate individual student reports"""
    individual_reports = {}
    
    for student_id in student_ids:
        student = db.query(StudentRecord).filter(
            StudentRecord.student_id == student_id,
            StudentRecord.user_id == user_id
        ).first()
        
        if student:
            # Get progress history
            history = db.query(ProgressHistory).filter(
                ProgressHistory.student_id == student_id
            ).order_by(ProgressHistory.recorded_at.desc()).all()
            
            # Calculate statistics
            recent_scores = [entry.percentage for entry in history[:10]]
            average_score = sum(recent_scores) / len(recent_scores) if recent_scores else 0
            
            individual_reports[student_id] = {
                "student": {
                    "name": student.name,
                    "class_name": student.class_name,
                    "grade": student.grade,
                    "progress": student.progress,
                    "status": student.status
                },
                "recent_performance": [
                    {
                        "assignment_name": entry.assignment_name,
                        "score": entry.score,
                        "max_score": entry.max_score,
                        "percentage": entry.percentage,
                        "date": entry.date_completed.isoformat()
                    }
                    for entry in history[:5]
                ],
                "statistics": {
                    "average_score": round(average_score, 2),
                    "total_assignments": len(history),
                    "trend": calculate_trend_db(history)
                }
            }
    
    return {
        "report_type": "individual_reports",
        "generated_at": datetime.now().isoformat(),
        "reports": individual_reports,
        "include_charts": include_charts
    }

async def generate_comparison_report(student_ids: List[str], include_charts: bool, user_id: int, db: Session):
    """Generate comparison report for multiple students"""
    comparison_data = {}
    
    for student_id in student_ids:
        student = db.query(StudentRecord).filter(
            StudentRecord.student_id == student_id,
            StudentRecord.user_id == user_id
        ).first()
        
        if student:
            # Get recent progress
            history = db.query(ProgressHistory).filter(
                ProgressHistory.student_id == student_id
            ).order_by(ProgressHistory.recorded_at.desc()).limit(5).all()
            
            comparison_data[student_id] = {
                "student_info": {
                    "name": student.name,
                    "class": student.class_name,
                    "grade": student.grade
                },
                "recent_performance": [
                    {
                        "assignment": entry.assignment_name,
                        "score": f"{entry.score}/{entry.max_score}",
                        "percentage": entry.percentage
                    }
                    for entry in history
                ],
                "statistics": {
                    "average_score": student.progress,
                    "attendance": student.attendance,
                    "status": student.status
                }
            }
    
    return {
        "report_type": "comparison_report",
        "generated_at": datetime.now().isoformat(),
        "comparison_data": comparison_data,
        "include_charts": include_charts
    }

def create_progress_report_pdf(report_data: Dict, filepath: str):
    """Create PDF progress report"""
    try:
        html_content = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                .header {{ text-align: center; border-bottom: 2px solid #333; padding-bottom: 20px; }}
                .summary {{ background: #f5f5f5; padding: 20px; margin: 20px 0; }}
                .student-table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
                .student-table th, .student-table td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
                .student-table th {{ background-color: #f2f2f2; }}
                .status-excellent {{ color: #00ff9d; font-weight: bold; }}
                .status-good {{ color: #64ffda; font-weight: bold; }}
                .status-average {{ color: #ffd166; font-weight: bold; }}
                .status-poor {{ color: #ff6b6b; font-weight: bold; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>Student Progress Report</h1>
                <p>Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>
            </div>
            
            <div class="summary">
                <h2>Class Summary</h2>
                <p><strong>Class:</strong> {report_data.get('class_name', 'All Classes')}</p>
                <p><strong>Total Students:</strong> {len(report_data.get('students', []))}</p>
                <p><strong>Average Progress:</strong> {report_data.get('summary', {}).get('average_progress', 0)}%</p>
            </div>
            
            <h2>Student Progress Details</h2>
            <table class="student-table">
                <thead>
                    <tr>
                        <th>Student Name</th>
                        <th>Class</th>
                        <th>Grade</th>
                        <th>Progress</th>
                        <th>Status</th>
                        <th>Last Updated</th>
                    </tr>
                </thead>
                <tbody>
                    {"".join([f"""
                    <tr>
                        <td>{student['name']}</td>
                        <td>{student['class_name']}</td>
                        <td>{student['grade']}</td>
                        <td>{student['progress']}%</td>
                        <td class="status-{student['status']}">{student['status'].title()}</td>
                        <td>{student['last_updated'][:10]}</td>
                    </tr>
                    """ for student in report_data.get('students', [])])}
                </tbody>
            </table>
        </body>
        </html>
        """
        
        create_pdf_from_html_optional(html_content, filepath)
        
    except Exception as e:
        print(f"Error creating PDF report: {e}")
        raise

def create_progress_report_docx(report_data: Dict, filepath: str):
    """Create DOCX progress report"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document()
        
        # Title
        doc.add_heading('Student Progress Report', 0)
        doc.add_paragraph(f'Generated on: {datetime.now().strftime("%Y-%m-%d %H:%M")}')
        
        # Summary
        doc.add_heading('Class Summary', level=1)
        summary = report_data.get('summary', {})
        doc.add_paragraph(f'Class: {report_data.get("class_name", "All Classes")}')
        doc.add_paragraph(f'Total Students: {len(report_data.get("students", []))}')
        doc.add_paragraph(f'Average Progress: {summary.get("average_progress", 0)}%')
        doc.add_paragraph(f'Attendance Rate: {summary.get("attendance_rate", 0)}%')
        
        # Student table
        doc.add_heading('Student Progress Details', level=1)
        table = doc.add_table(rows=1, cols=6)
        table.style = 'Light Grid Accent 1'
        
        # Header row
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Student Name'
        hdr_cells[1].text = 'Class'
        hdr_cells[2].text = 'Grade'
        hdr_cells[3].text = 'Progress'
        hdr_cells[4].text = 'Status'
        hdr_cells[5].text = 'Last Updated'
        
        # Data rows
        for student in report_data.get('students', []):
            row_cells = table.add_row().cells
            row_cells[0].text = student['name']
            row_cells[1].text = student['class_name']
            row_cells[2].text = student['grade']
            row_cells[3].text = f"{student['progress']}%"
            row_cells[4].text = student['status'].title()
            row_cells[5].text = student['last_updated'][:10]
        
        doc.save(filepath)
        
    except Exception as e:
        print(f"Error creating DOCX report: {e}")
        # Fallback to PDF if DOCX fails
        create_progress_report_pdf(report_data, filepath.replace('.docx', '.pdf'))

@app.get("/api/analytics/performance-trends")
async def get_performance_trends(
    class_name: str = None,
    weeks: int = 8,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Get performance trends data for charts
    """
    try:
        # Get students
        query = db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id,
            StudentRecord.is_active == True
        )
        
        if class_name and class_name != "all":
            query = query.filter(StudentRecord.class_name == class_name)
        
        students = query.all()
        
        trends_data = generate_performance_trends_db(students, f"last_{weeks}_weeks", db)
        
        return JSONResponse(content={
            "trends": trends_data,
            "time_period": f"Last {weeks} weeks",
            "class": class_name or "All Classes"
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating trends: {str(e)}")

# ---------- Download with Database Tracking ----------
@app.get("/download/{filename}")
async def download(
    filename: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    
    if not os.path.exists(file_path):
        # Check if file exists in database
        generated_file = db.query(GeneratedFile).filter(
            GeneratedFile.filename == filename,
            GeneratedFile.user_id == current_user.id,
            GeneratedFile.is_active == True
        ).first()
        
        if not generated_file:
            raise HTTPException(status_code=404, detail="File not found.")
        
        # Check if file has expired
        if generated_file.expires_at and generated_file.expires_at < datetime.utcnow():
            generated_file.is_active = False
            db.commit()
            raise HTTPException(status_code=410, detail="File has expired.")
        
        file_path = generated_file.file_path
    
    # Determine media type
    if filename.lower().endswith('.pdf'):
        media_type = "application/pdf"
    elif filename.lower().endswith('.docx'):
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif filename.lower().endswith('.pptx'):
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif filename.lower().endswith('.json'):
        media_type = "application/json"
    elif filename.lower().endswith('.zip'):
        media_type = "application/zip"
    else:
        media_type = "application/octet-stream"
    
    return FileResponse(
        file_path, 
        media_type=media_type, 
        filename=filename,
        headers={'Content-Disposition': f'attachment; filename="{filename}"'}
    )

# ---------- User Data Management ----------
@app.get("/api/user/data")
async def get_user_data(
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Get all user data (conversations, lesson plans, files, students)"""
    try:
        # Get conversations
        conversations = db.query(Conversation).filter(
            Conversation.user_id == current_user.id,
            Conversation.is_active == True
        ).order_by(Conversation.updated_at.desc()).limit(50).all()
        
        # Get lesson plans
        lesson_plans = db.query(LessonPlan).filter(
            LessonPlan.user_id == current_user.id,
            LessonPlan.is_active == True
        ).order_by(LessonPlan.updated_at.desc()).limit(50).all()
        
        # Get generated files
        generated_files = db.query(GeneratedFile).filter(
            GeneratedFile.user_id == current_user.id,
            GeneratedFile.is_active == True
        ).order_by(GeneratedFile.created_at.desc()).limit(50).all()
        
        # Get students
        students = db.query(StudentRecord).filter(
            StudentRecord.user_id == current_user.id,
            StudentRecord.is_active == True
        ).order_by(StudentRecord.last_updated.desc()).all()
        
        return JSONResponse(content={
            "status": "success",
            "user": {
                "id": current_user.id,
                "username": current_user.username,
                "email": current_user.email,
                "full_name": current_user.full_name,
                "account_type": current_user.account_type,
                "created_at": current_user.created_at.isoformat(),
                "last_login": current_user.last_login.isoformat() if current_user.last_login else None
            },
            "stats": {
                "conversations": len(conversations),
                "lesson_plans": len(lesson_plans),
                "generated_files": len(generated_files),
                "students": len(students)
            },
            "data": {
                "conversations": [
                    {
                        "id": conv.conversation_id,
                        "title": conv.title,
                        "created_at": conv.created_at.isoformat(),
                        "updated_at": conv.updated_at.isoformat()
                    }
                    for conv in conversations
                ],
                "lesson_plans": [
                    {
                        "id": lp.lesson_id,
                        "title": lp.title,
                        "topic": lp.topic,
                        "grade": lp.grade,
                        "created_at": lp.created_at.isoformat()
                    }
                    for lp in lesson_plans
                ],
                "generated_files": [
                    {
                        "id": gf.file_id,
                        "filename": gf.filename,
                        "file_type": gf.file_type,
                        "created_at": gf.created_at.isoformat(),
                        "expires_at": gf.expires_at.isoformat() if gf.expires_at else None
                    }
                    for gf in generated_files
                ],
                "students": [
                    {
                        "id": student.student_id,
                        "name": student.name,
                        "class": student.class_name,
                        "grade": student.grade,
                        "progress": student.progress,
                        "status": student.status
                    }
                    for student in students
                ]
            }
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching user data: {str(e)}")

# ---------- Conversations ----------
@app.get("/conversation/{conversation_id}")
async def get_conversation(
    conversation_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Get conversation with messages"""
    try:
        conversation = db.query(Conversation).filter(
            Conversation.conversation_id == conversation_id,
            Conversation.user_id == current_user.id
        ).first()
        
        if not conversation:
            raise HTTPException(status_code=404, detail="Conversation not found")
        
        # Get messages
        messages = db.query(ConversationMessage).filter(
            ConversationMessage.conversation_id == conversation_id,
            ConversationMessage.user_id == current_user.id
        ).order_by(ConversationMessage.created_at).all()
        
        return JSONResponse(content={
            "id": conversation.conversation_id,
            "title": conversation.title,
            "created_at": conversation.created_at.isoformat(),
            "updated_at": conversation.updated_at.isoformat(),
            "messages": [
                {
                    "role": msg.role,
                    "content": msg.content,
                    "message_type": msg.message_type,
                    "created_at": msg.created_at.isoformat()
                }
                for msg in messages
            ]
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching conversation: {str(e)}")

@app.get("/conversations/")
async def list_conversations(
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """List all conversations for user"""
    try:
        conversations = db.query(Conversation).filter(
            Conversation.user_id == current_user.id,
            Conversation.is_active == True
        ).order_by(Conversation.updated_at.desc()).all()
        
        return JSONResponse(content=[
            {
                "id": conv.conversation_id,
                "title": conv.title,
                "created_at": conv.created_at.isoformat(),
                "updated_at": conv.updated_at.isoformat(),
                "message_count": db.query(ConversationMessage).filter(
                    ConversationMessage.conversation_id == conv.conversation_id
                ).count()
            }
            for conv in conversations
        ])
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error listing conversations: {str(e)}")

@app.delete("/conversation/{conversation_id}")
async def delete_conversation(
    conversation_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Delete conversation"""
    try:
        conversation = db.query(Conversation).filter(
            Conversation.conversation_id == conversation_id,
            Conversation.user_id == current_user.id
        ).first()
        
        if not conversation:
            raise HTTPException(status_code=404, detail="Conversation not found")
        
        # Soft delete (mark as inactive)
        conversation.is_active = False
        db.commit()
        
        return {"message": "Conversation deleted successfully"}
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Error deleting conversation: {str(e)}")

# ---------- Cleanup Old Files from Database ----------
def cleanup_old_files_db():
    """Clean up expired files from database and filesystem"""
    try:
        db = SessionLocal()
        current_time = datetime.utcnow()
        
        # Find expired files
        expired_files = db.query(GeneratedFile).filter(
            GeneratedFile.is_active == True,
            GeneratedFile.expires_at < current_time
        ).all()
        
        for file in expired_files:
            # Delete from filesystem
            file_path = file.file_path
            if os.path.exists(file_path):
                os.remove(file_path)
            
            # Mark as inactive in database
            file.is_active = False
        
        # Commit changes
        db.commit()
        
        # Also clean up orphaned files in OUTPUT_FOLDER older than 24 hours
        for filename in os.listdir(OUTPUT_FOLDER):
            file_path = os.path.join(OUTPUT_FOLDER, filename)
            if os.path.isfile(file_path):
                file_age = current_time.timestamp() - os.path.getctime(file_path)
                if file_age > 24 * 3600:  # 24 hours
                    # Check if file exists in database
                    db_file = db.query(GeneratedFile).filter(
                        GeneratedFile.filename == filename
                    ).first()
                    
                    if not db_file:
                        os.remove(file_path)
                        print(f"Cleaned up orphaned file: {filename}")
        
        db.close()
        print(f"Cleaned up {len(expired_files)} expired files from database")
        
    except Exception as e:
        print(f"Database cleanup error: {e}")

# ---------- Health Check with Database ----------
@app.get("/health")
async def health_check():
    db_healthy = test_connection()
    return {
        "status": "healthy" if db_healthy else "degraded",
        "database": "connected" if db_healthy else "disconnected",
        "timestamp": datetime.now().isoformat(),
        "service": "Sahayak AI Teaching Assistant"
    }

# ---------- Database Maintenance ----------
@app.post("/api/admin/cleanup")
async def admin_cleanup(
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Admin endpoint to clean up old data"""
    # Check if user is admin
    if current_user.account_type != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    
    try:
        # Clean up old files
        cleanup_old_files_db()
        
        # Clean up old conversations (older than 30 days)
        cutoff_date = datetime.utcnow() - timedelta(days=30)
        old_conversations = db.query(Conversation).filter(
            Conversation.updated_at < cutoff_date,
            Conversation.is_active == True
        ).all()
        
        for conv in old_conversations:
            conv.is_active = False
        
        # Clean up old lesson plans (older than 30 days)
        old_lesson_plans = db.query(LessonPlan).filter(
            LessonPlan.updated_at < cutoff_date,
            LessonPlan.is_active == True
        ).all()
        
        for lp in old_lesson_plans:
            lp.is_active = False
        
        db.commit()
        
        return JSONResponse(content={
            "status": "success",
            "message": f"Cleaned up {len(old_conversations)} old conversations and {len(old_lesson_plans)} old lesson plans",
            "timestamp": datetime.now().isoformat()
        })
        
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Cleanup failed: {str(e)}")

# ---------- Startup ----------
@app.on_event("startup")
async def startup_event():
    print("🚀 Sahayak Backend started successfully!")
    
    # Test database connection
    if test_connection():
        print("✅ Database connected successfully")
    else:
        print("⚠️  Database connection failed - some features may not work")
    
    # Initial cleanup
    cleanup_old_files_db()
    
    # Schedule cleanup to run daily at 3 AM
    schedule.every().day.at("03:00").do(cleanup_old_files_db)
    
    # Start cleanup thread
    cleanup_thread = threading.Thread(target=run_scheduler, daemon=True)
    cleanup_thread.start()

def run_scheduler():
    while True:
        schedule.run_pending()
        time.sleep(3600)  # Check every hour

# ---------- Errors ----------
@app.exception_handler(404)
async def not_found_handler(request, exc):
    return JSONResponse(status_code=404, content={"error": "Endpoint not found", "detail": str(exc)})

@app.exception_handler(500)
async def internal_error_handler(request, exc):
    return JSONResponse(status_code=500, content={"error": "Internal server error", "detail": str(exc)})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=int(os.environ.get("PORT", 8000)), reload=True)